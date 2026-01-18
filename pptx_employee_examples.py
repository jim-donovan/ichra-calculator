"""
Employee Examples Slide Generator

Generates PowerPoint slides for representative employee examples
(youngest employee, mid-age family, oldest employee) with cost comparison tables.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from typing import Dict, List, Optional
from io import BytesIO
from pathlib import Path
from dataclasses import dataclass, field
from datetime import datetime, timedelta


# Color scheme matching ICHRA dashboard design
COLORS = {
    # Current/Renewal headers
    'current_bg': RGBColor(0xEF, 0xF6, 0xFF),       # #EFF6FF - blue-50
    'current_text': RGBColor(0x1C, 0x39, 0x8E),     # #1C398E
    'renewal_bg': RGBColor(0xFE, 0xF2, 0xF2),       # #FEF2F2 - red-50
    'renewal_text': RGBColor(0x82, 0x18, 0x1A),     # #82181A

    # ICHRA Metal headers
    'bronze_bg': RGBColor(0xFE, 0xF3, 0xC7),        # #FEF3C7 - amber-100
    'bronze_text': RGBColor(0x92, 0x40, 0x0E),      # #92400E
    'silver_bg': RGBColor(0xF3, 0xF4, 0xF6),        # #F3F4F6 - gray-100
    'silver_text': RGBColor(0x37, 0x41, 0x51),      # #374151
    'gold_bg': RGBColor(0xFE, 0xF9, 0xC3),          # #FEF9C3 - yellow-100
    'gold_text': RGBColor(0x85, 0x4D, 0x0E),        # #854D0E

    # Cooperative (HAS)
    'coop_bg': RGBColor(0xF0, 0xFD, 0xF4),          # #F0FDF4 - green-50
    'coop_text': RGBColor(0x0D, 0x54, 0x2B),        # #0D542B

    # Sedera
    'sedera_bg': RGBColor(0xFE, 0xF3, 0xC7),        # #FEF3C7 - amber-100
    'sedera_text': RGBColor(0x92, 0x40, 0x0E),      # #92400E

    # Neutral
    'row_label_bg': RGBColor(0xF9, 0xFA, 0xFB),     # #F9FAFB
    'row_label_text': RGBColor(0x36, 0x41, 0x53),   # #364153
    'total_bg': RGBColor(0xF3, 0xF4, 0xF6),         # #F3F4F6
    'border': RGBColor(0xE5, 0xE7, 0xEB),           # #E5E7EB
    'white': RGBColor(0xFF, 0xFF, 0xFF),            # #FFFFFF
    'black': RGBColor(0x10, 0x18, 0x28),            # #101828
    'secondary': RGBColor(0x4A, 0x55, 0x65),        # #4A5565

    # Insight box
    'insight_bg': RGBColor(0xEF, 0xF6, 0xFF),       # #EFF6FF
    'insight_border': RGBColor(0xBE, 0xDB, 0xFF),   # #BEDBFF
    'insight_text': RGBColor(0x1C, 0x39, 0x8E),     # #1C398E

    # Savings/cost
    'savings_green': RGBColor(0x00, 0xA6, 0x3E),    # #00A63E
    'cost_red': RGBColor(0xDC, 0x26, 0x26),         # #DC2626
}

# Slide dimensions (standard widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Banner image path
BANNER_IMAGE = Path(__file__).parent / "decoratives" / "PPT_header.png"


@dataclass
class EmployeeExampleData:
    """Data for a single employee example slide"""
    label: str  # "Youngest Employee", "Mid-Age Family", "Older Employee"
    name: str
    age: int
    tier: str  # "Employee Only", "Family", etc.
    location: str
    family_ages: List[Dict] = field(default_factory=list)
    family_status: str = "EE"  # EE, ES, EC, F

    # Cost breakdown by scenario
    costs: Dict[str, Dict[str, float]] = field(default_factory=dict)
    # e.g., costs['Current'] = {'employee': 250, 'employer': 500}

    winner: str = ""  # Which option is best
    insight: str = ""  # Key insight text

    # Plan details for the expanded view
    metal_plan_details: Dict = field(default_factory=dict)
    member_breakdowns: Dict = field(default_factory=dict)  # Individual member rates by plan
    current_total_monthly: float = 0.0
    renewal_total_monthly: float = 0.0
    use_ee_rate_only: bool = False
    has_renewal_data: bool = True  # When False, hide renewal column and compare vs current
    contribution_strategy: str = ""  # Description of employer contribution strategy


class EmployeeExamplesSlideGenerator:
    """Generate PowerPoint slides for employee examples"""

    def __init__(self, client_name: str = "", plan_config: Dict = None, include_qr_links: bool = False):
        """Initialize with a new blank presentation

        Args:
            client_name: Optional client name for footer
            plan_config: Plan configurator settings dict with 'hap_enabled', 'hap_iuas',
                        'sedera_enabled', 'sedera_iuas' keys
            include_qr_links: Whether to add QR codes linking to member breakdown pages
        """
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.client_name = client_name
        self.plan_config = plan_config or {}
        self.include_qr_links = include_qr_links

    def _set_cell_fill(self, cell, color: RGBColor):
        """Set cell background color"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    def _set_cell_text(self, cell, text: str, color: RGBColor,
                       font_size: int = 11, bold: bool = False,
                       align: PP_ALIGN = PP_ALIGN.CENTER):
        """Set cell text with formatting"""
        cell.text = str(text) if text is not None else "—"
        cell.text_frame.paragraphs[0].alignment = align
        cell.text_frame.paragraphs[0].font.name = 'Poppins'
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.text_frame.paragraphs[0].font.color.rgb = color
        cell.text_frame.paragraphs[0].font.bold = bold
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _add_cell_line(self, cell, text: str, color: RGBColor,
                       font_size: int = 10, bold: bool = False,
                       align: PP_ALIGN = PP_ALIGN.CENTER):
        """Add additional line to cell"""
        p = cell.text_frame.add_paragraph()
        p.text = str(text) if text is not None else "—"
        p.alignment = align
        p.font.name = 'Poppins'
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold

    def _format_currency(self, val) -> str:
        """Format value as currency"""
        if val is None:
            return "—"
        return f"${val:,.0f}"

    def _add_banner(self, slide) -> None:
        """Add banner image at the top of the slide"""
        if BANNER_IMAGE.exists():
            # Add banner spanning full width at the top
            slide.shapes.add_picture(
                str(BANNER_IMAGE),
                left=Inches(0),
                top=Inches(0),
                width=SLIDE_WIDTH,
                height=Inches(0.25)  # Compact banner height
            )

    def _build_family_string(self, family_ages: List[Dict]) -> str:
        """Build family ages string"""
        if not family_ages:
            return ""

        age_parts = []
        spouse_ages = [f"Spouse ({fa['age']})" for fa in family_ages if fa.get('relationship', '').lower() == 'spouse']
        # Get child ages and sort in ascending order
        child_ages = sorted([fa['age'] for fa in family_ages if fa.get('relationship', '').lower() == 'child'])

        if spouse_ages:
            age_parts.extend(spouse_ages)
        if child_ages:
            if len(child_ages) == 1:
                age_parts.append(f"Child ({child_ages[0]})")
            else:
                age_parts.append(f"Children ({', '.join(str(a) for a in child_ages)})")

        return " | ".join(age_parts) if age_parts else ""

    def add_employee_slide(self, employee: EmployeeExampleData):
        """Add a slide for a single employee example"""
        # Add blank slide
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add banner at top
        self._add_banner(slide)

        # Vertical offset to account for banner (compressed layout)
        banner_offset = Inches(0.25)

        # --- COMPANY NAME (top right) ---
        if self.client_name:
            company_box = slide.shapes.add_textbox(
                Inches(9.5), Inches(0.35), Inches(3.5), Inches(0.3)
            )
            company_tf = company_box.text_frame
            p = company_tf.paragraphs[0]
            p.text = self.client_name
            p.font.name = 'Poppins'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLORS['black']
            p.alignment = PP_ALIGN.RIGHT

        # --- HEADER SECTION ---
        # Title: "Youngest Employee: John Smith"
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.35) + banner_offset, Inches(9), Inches(0.5)
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p = title_tf.paragraphs[0]
        p.text = f"{employee.label}: {employee.name}"
        p.font.name = 'Poppins'
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLORS['black']

        # Subtitle: Age, Tier, Family, Location (combined)
        family_str = self._build_family_string(employee.family_ages)
        subtitle_text = f"Age {employee.age} | {employee.tier} | {employee.location}"
        if family_str:
            subtitle_text = f"Age {employee.age} | {employee.tier} | {family_str} | {employee.location}"

        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.75) + banner_offset, Inches(12), Inches(0.3)
        )
        subtitle_tf = subtitle_box.text_frame
        p = subtitle_tf.paragraphs[0]
        p.text = subtitle_text
        p.font.name = 'Poppins'
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['secondary']

        # Contribution strategy (below subtitle)
        if employee.contribution_strategy:
            strategy_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.0) + banner_offset, Inches(12), Inches(0.3)
            )
            strategy_tf = strategy_box.text_frame
            p = strategy_tf.paragraphs[0]
            p.text = employee.contribution_strategy
            p.font.name = 'Poppins'
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['secondary']

        # --- COST COMPARISON TABLE ---
        # Build dynamic column list based on plan_config and has_renewal_data

        # Base columns: Label, Current, (Renewal if available), Bronze, Silver, Gold
        columns = [
            {'header': '', 'scenario': None, 'bg': COLORS['white'], 'text': COLORS['black']},
            {'header': 'Current', 'scenario': 'Current', 'bg': COLORS['current_bg'], 'text': COLORS['current_text']},
        ]
        # Only include Renewal column if renewal data is available
        if employee.has_renewal_data:
            columns.append({'header': 'Renewal', 'scenario': 'Renewal', 'bg': COLORS['renewal_bg'], 'text': COLORS['renewal_text']})

        # Metal plan columns
        columns.extend([
            {'header': 'Lowest Cost\nBronze', 'scenario': 'ICHRA Bronze', 'bg': COLORS['bronze_bg'], 'text': COLORS['bronze_text']},
            {'header': 'Lowest Cost\nSilver', 'scenario': 'ICHRA Silver', 'bg': COLORS['silver_bg'], 'text': COLORS['silver_text']},
            {'header': 'Lowest Cost\nGold', 'scenario': 'ICHRA Gold', 'bg': COLORS['gold_bg'], 'text': COLORS['gold_text']},
        ])

        # Add HAS columns for each enabled IUA level
        if self.plan_config.get('hap_enabled') and self.plan_config.get('hap_iuas'):
            sorted_hap_iuas = sorted(self.plan_config['hap_iuas'], key=lambda x: float(x.replace('k', '')))
            for iua in sorted_hap_iuas:
                key = f"HAS ${iua}"
                if key in employee.costs:
                    columns.append({
                        'header': f'HAS ${iua}', 'scenario': key,
                        'bg': COLORS['coop_bg'], 'text': COLORS['coop_text']
                    })

        # Add Sedera columns for each enabled IUA level
        if self.plan_config.get('sedera_enabled') and self.plan_config.get('sedera_iuas'):
            sorted_sedera_iuas = sorted(self.plan_config['sedera_iuas'], key=lambda x: int(x))
            for iua in sorted_sedera_iuas:
                # Match the display format used in build_employee_example
                iua_display = iua if int(iua) < 1000 else f"{int(iua)//1000}k" if int(iua) % 1000 == 0 else f"{int(iua)/1000}k"
                key = f"Sedera ${iua_display}"
                if key in employee.costs:
                    columns.append({
                        'header': f'Sedera ${iua_display}', 'scenario': key,
                        'bg': COLORS['sedera_bg'], 'text': COLORS['sedera_text']
                    })

        # Table dimensions - dynamically calculate column widths to fit
        table_left = Inches(0.5)
        table_top = Inches(1.25) + banner_offset
        num_cols = len(columns)
        max_table_width = 12.333  # inches - matches slide content width

        # Label column gets fixed smaller width, data columns share remaining space
        label_col_width = 1.0  # inches
        data_col_width = (max_table_width - label_col_width) / (num_cols - 1) if num_cols > 1 else 1.5

        # Build column widths list
        col_widths = [Inches(label_col_width)] + [Inches(data_col_width)] * (num_cols - 1)
        table_width = max_table_width

        # Create table (4 rows: header, employee, employer, total)
        table = slide.shapes.add_table(
            rows=4, cols=num_cols,
            left=table_left, top=table_top,
            width=Inches(table_width), height=Inches(1.8)
        ).table

        # Set column widths
        for i, width in enumerate(col_widths):
            table.columns[i].width = width

        # Row 0: Headers (handle multi-line headers)
        for col, col_def in enumerate(columns):
            cell = table.cell(0, col)
            self._set_cell_fill(cell, col_def['bg'])
            header_text = col_def['header']
            if '\n' in header_text:
                lines = header_text.split('\n')
                self._set_cell_text(cell, lines[0], col_def['text'], font_size=10, bold=True)
                self._add_cell_line(cell, lines[1], col_def['text'], font_size=11, bold=True)
            else:
                self._set_cell_text(cell, header_text, col_def['text'], font_size=11, bold=True)

        # Get costs data
        costs = employee.costs

        # Row 1: Employee costs
        for col, col_def in enumerate(columns):
            cell = table.cell(1, col)
            if col == 0:
                self._set_cell_fill(cell, COLORS['row_label_bg'])
                self._set_cell_text(cell, "Employee", COLORS['row_label_text'], font_size=11, bold=True, align=PP_ALIGN.LEFT)
            else:
                self._set_cell_fill(cell, COLORS['white'])
                val = costs.get(col_def['scenario'], {}).get('employee', 0)
                self._set_cell_text(cell, self._format_currency(val), COLORS['black'], font_size=11)

        # Row 2: Employer costs
        for col, col_def in enumerate(columns):
            cell = table.cell(2, col)
            if col == 0:
                self._set_cell_fill(cell, COLORS['row_label_bg'])
                self._set_cell_text(cell, "Employer", COLORS['row_label_text'], font_size=11, bold=True, align=PP_ALIGN.LEFT)
            else:
                self._set_cell_fill(cell, COLORS['white'])
                val = costs.get(col_def['scenario'], {}).get('employer', 0)
                self._set_cell_text(cell, self._format_currency(val), COLORS['black'], font_size=11)

        # Row 3: Total costs
        for col, col_def in enumerate(columns):
            cell = table.cell(3, col)
            self._set_cell_fill(cell, COLORS['total_bg'])
            if col == 0:
                self._set_cell_text(cell, "Total", COLORS['black'], font_size=11, bold=True, align=PP_ALIGN.LEFT)
            else:
                ee = costs.get(col_def['scenario'], {}).get('employee', 0)
                er = costs.get(col_def['scenario'], {}).get('employer', 0)
                self._set_cell_text(cell, self._format_currency(ee + er), COLORS['black'], font_size=11, bold=True)

        # --- INSIGHT BOX ---
        insight_top = Inches(3.3) + banner_offset
        insight_box = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), insight_top,
            Inches(12.333), Inches(0.6)
        )
        insight_box.fill.solid()
        insight_box.fill.fore_color.rgb = COLORS['insight_bg']
        insight_box.line.color.rgb = COLORS['insight_border']
        insight_box.line.width = Pt(1)

        insight_text_box = slide.shapes.add_textbox(
            Inches(0.7), insight_top + Inches(0.15),
            Inches(11.9), Inches(0.4)
        )
        insight_tf = insight_text_box.text_frame
        p = insight_tf.paragraphs[0]
        p.text = f"💡 {employee.insight}"
        p.font.name = 'Poppins'
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['insight_text']

        # --- PLAN DETAILS TABLE ---
        import logging
        logger = logging.getLogger(__name__)
        has_dependents = employee.family_status in ('ES', 'EC', 'F') or bool(employee.family_ages)
        will_add_qr = self.include_qr_links and bool(employee.member_breakdowns) and has_dependents
        logger.info(f"Employee {employee.name}: family_status={employee.family_status}, has_dependents={has_dependents}, will_add_qr={will_add_qr}")

        if employee.metal_plan_details:
            details_label = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.1) + banner_offset, Inches(3), Inches(0.3)
            )
            details_tf = details_label.text_frame
            p = details_tf.paragraphs[0]
            p.text = "Plan Details"
            p.font.name = 'Poppins'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLORS['black']

            self._add_plan_details_table(slide, employee, banner_offset)

        # --- FOOTER ---
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.0), Inches(12), Inches(0.3)
        )
        footer_tf = footer_box.text_frame
        p = footer_tf.paragraphs[0]
        date_str = datetime.now().strftime("%m.%d.%y")
        if self.client_name:
            footer_text = f"Generated by Glove Benefits for {self.client_name} | {date_str}"
        else:
            footer_text = f"Generated by Glove Benefits | {date_str}"
        p.text = footer_text
        p.font.name = 'Poppins'
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['secondary']

    def _add_plan_details_table(self, slide, employee: EmployeeExampleData, banner_offset=Inches(0)):
        """Add the plan details comparison table"""
        details = employee.metal_plan_details

        # Extract plan data for each metal
        metals = ['Bronze', 'Silver', 'Gold']
        available_metals = [m for m in metals if m in details and details[m]]

        if not available_metals:
            return

        # Table dimensions
        table_left = Inches(0.5)
        table_top = Inches(4.45) + banner_offset
        num_cols = len(available_metals) + 1  # +1 for row labels
        col_width = Inches(2.5)
        table_width = col_width * num_cols

        # Rows: Plan, Premium, vs Renewal, Deductible, MOOP
        num_rows = 6
        table = slide.shapes.add_table(
            rows=num_rows, cols=num_cols,
            left=table_left, top=table_top,
            width=table_width, height=Inches(1.8)
        ).table

        # Set column widths
        for i in range(num_cols):
            table.columns[i].width = col_width

        # Header row
        header_cell = table.cell(0, 0)
        self._set_cell_fill(header_cell, COLORS['white'])
        self._set_cell_text(header_cell, "", COLORS['black'])

        metal_colors = {
            'Bronze': (COLORS['bronze_bg'], COLORS['bronze_text']),
            'Silver': (COLORS['silver_bg'], COLORS['silver_text']),
            'Gold': (COLORS['gold_bg'], COLORS['gold_text']),
        }

        for col, metal in enumerate(available_metals, 1):
            cell = table.cell(0, col)
            bg, text = metal_colors.get(metal, (COLORS['white'], COLORS['black']))
            self._set_cell_fill(cell, bg)
            # Two-line header: "Lowest Cost" on first line, metal name on second
            self._set_cell_text(cell, "Lowest Cost", text, font_size=10, bold=True)
            self._add_cell_line(cell, metal, text, font_size=11, bold=True)

        # Data rows - use dynamic comparison label based on renewal data availability
        comparison_label = "vs Renewal (Savings)" if employee.has_renewal_data else "vs Current (Savings)"
        # Reordered: Plan Name, Deductible, Max OOP, then Premium and comparison
        row_labels = ["Plan Name", "Deductible", "Max Out-of-Pocket", "Monthly Premium", comparison_label]
        # Use renewal or current as comparison baseline
        comparison_baseline = employee.renewal_total_monthly if employee.has_renewal_data else employee.current_total_monthly

        for row_idx, label in enumerate(row_labels, 1):
            # Row label
            label_cell = table.cell(row_idx, 0)
            self._set_cell_fill(label_cell, COLORS['row_label_bg'])
            self._set_cell_text(label_cell, label, COLORS['row_label_text'], font_size=10, bold=True, align=PP_ALIGN.LEFT)

            # Values for each metal
            for col, metal in enumerate(available_metals, 1):
                cell = table.cell(row_idx, col)
                self._set_cell_fill(cell, COLORS['white'])

                metal_data = details.get(metal, {})

                if label == "Plan Name":
                    plan_name = metal_data.get('plan_name', '—')
                    # Enable word wrap for long plan names
                    cell.text_frame.word_wrap = True
                    self._set_cell_text(cell, plan_name, COLORS['black'], font_size=9)

                elif label == "Deductible":
                    ded = metal_data.get('deductible')
                    self._set_cell_text(cell, self._format_currency(ded), COLORS['black'], font_size=10)

                elif label == "Max Out-of-Pocket":
                    moop = metal_data.get('moop')
                    self._set_cell_text(cell, self._format_currency(moop), COLORS['black'], font_size=10)

                elif label == "Monthly Premium":
                    premium = metal_data.get('premium', 0)
                    self._set_cell_text(cell, self._format_currency(premium), COLORS['black'], font_size=10)

                elif label == comparison_label:  # "vs Renewal" or "vs Current"
                    premium = metal_data.get('premium', 0)
                    if comparison_baseline and comparison_baseline > 0:
                        diff = premium - comparison_baseline
                        pct = abs(diff) / comparison_baseline * 100
                        if diff < 0:
                            # Savings: show without minus sign, add percentage
                            text = f"${abs(diff):,.0f} ({pct:.0f}%)"
                            color = COLORS['savings_green']
                        elif diff > 0:
                            # Cost increase: show with plus sign, add percentage
                            text = f"+${diff:,.0f} ({pct:.0f}%)"
                            color = COLORS['cost_red']
                        else:
                            text = "$0"
                            color = COLORS['secondary']
                        self._set_cell_text(cell, text, color, font_size=10, bold=True)
                    else:
                        self._set_cell_text(cell, "—", COLORS['secondary'], font_size=10)

        # Add QR code linking to member breakdown page (if enabled, has breakdown data, and has dependents)
        has_dependents = employee.family_status in ('ES', 'EC', 'F') or bool(employee.family_ages)
        if self.include_qr_links and employee.member_breakdowns and has_dependents:
            self._add_qr_code(slide, employee, banner_offset)

    def _add_qr_code(self, slide, employee: EmployeeExampleData, banner_offset=Inches(0)):
        """Add QR code linking to member rate breakdown page."""
        import logging
        logger = logging.getLogger(__name__)
        logger.info(f"Attempting to add QR code for {employee.name}")

        try:
            from r2_storage import R2StorageService
            from qr_generator import generate_qr_code
            from member_breakdown_template import generate_member_breakdown_html
            from url_shortener import shorten_url

            # Check if R2 is configured
            r2 = R2StorageService()
            is_configured, error = r2.is_configured()
            if not is_configured:
                # Skip QR code if R2 not configured (graceful degradation)
                logger.warning(f"R2 not configured: {error}")
                return

            # Generate HTML page
            html = generate_member_breakdown_html(
                employee_name=employee.name,
                employee_age=employee.age,
                tier=employee.tier,
                location=employee.location,
                family_ages=employee.family_ages,
                member_breakdowns=employee.member_breakdowns,
                client_name=self.client_name
            )

            # Upload to R2 and get presigned URL
            logger.info(f"Uploading HTML to R2 for {employee.name}")
            url = r2.upload_html(html)
            if not url:
                logger.warning(f"Failed to upload HTML to R2 for {employee.name}")
                return
            logger.info(f"Uploaded successfully, URL: {url[:80]}...")

            # Shorten URL for cleaner QR code (falls back to original if dub not configured)
            short_url = shorten_url(url)
            if short_url != url:
                logger.info(f"Shortened URL: {short_url}")

            # Generate QR code image
            qr_buffer = generate_qr_code(short_url)
            if not qr_buffer:
                logger.warning(f"Failed to generate QR code for {employee.name}")
                return
            logger.info(f"QR code generated for {employee.name}")

            # Add QR code to slide (bottom-right corner)
            qr_size = Inches(2.5)
            qr_left = Inches(10.6)
            qr_top = Inches(4.25) + banner_offset

            slide.shapes.add_picture(
                qr_buffer,
                left=qr_left,
                top=qr_top,
                width=qr_size,
                height=qr_size
            )

            # Calculate expiration date (7 days from now)
            expiry_date = (datetime.now() + timedelta(days=7)).strftime("%b %d, %Y")

            # Add styled label below QR
            label_left = Inches(10.62)
            label_top = Inches(6.81)
            qr_label = slide.shapes.add_textbox(
                label_left, label_top, Inches(2.5), Inches(0.5)
            )
            tf = qr_label.text_frame
            tf.word_wrap = True

            # Title line
            p1 = tf.paragraphs[0]
            p1.text = "Scan for member details"
            p1.font.name = 'Poppins'
            p1.font.size = Pt(9)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(0x93, 0x42, 0x1B)  # Glove brand color
            p1.alignment = PP_ALIGN.CENTER

            # Clickable link line (for desktop viewers)
            p2 = tf.add_paragraph()
            # "click " text (no hyperlink)
            run1 = p2.add_run()
            run1.text = "click "
            run1.font.name = 'Poppins'
            run1.font.size = Pt(8)
            run1.font.color.rgb = COLORS['secondary']
            # "here" as hyperlink
            run2 = p2.add_run()
            run2.text = "here"
            run2.font.name = 'Poppins'
            run2.font.size = Pt(8)
            run2.font.color.rgb = RGBColor(0x37, 0xBE, 0xAE)  # Glove teal
            run2.font.underline = True
            run2.hyperlink.address = short_url
            p2.alignment = PP_ALIGN.CENTER

            # Expiry line
            p3 = tf.add_paragraph()
            p3.text = f"Expires {expiry_date}"
            p3.font.name = 'Poppins'
            p3.font.size = Pt(7)
            p3.font.color.rgb = COLORS['secondary']
            p3.alignment = PP_ALIGN.CENTER

        except ImportError as e:
            # Dependencies not installed, skip QR code
            pass
        except Exception as e:
            # Log error but don't fail the PPT generation
            import logging
            logging.warning(f"Failed to add QR code: {e}")

    def generate(self, employees: List[EmployeeExampleData]) -> BytesIO:
        """Generate the presentation with all employee slides"""
        for employee in employees:
            self.add_employee_slide(employee)

        # Save to BytesIO
        output = BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output


def generate_employee_examples_pptx(
    employee_examples: List[Dict],
    client_name: str = "",
    plan_config: Dict = None,
    include_qr_links: bool = False,
    has_renewal_data: bool = True
) -> BytesIO:
    """
    Generate PowerPoint slides for employee examples.

    Args:
        employee_examples: List of employee example dicts from DashboardData
        client_name: Optional client name for footer
        plan_config: Optional plan configurator settings dict for dynamic columns
        include_qr_links: Whether to add QR codes linking to member breakdown pages
        has_renewal_data: Whether renewal data is available (hides Renewal column if False)

    Returns:
        BytesIO containing the PowerPoint file
    """
    # Convert dicts to dataclass instances
    employees = []
    for emp in employee_examples:
        employees.append(EmployeeExampleData(
            label=emp.get('label', ''),
            name=emp.get('name', ''),
            age=emp.get('age', 0),
            tier=emp.get('tier', ''),
            location=emp.get('location', ''),
            family_ages=emp.get('family_ages', []),
            family_status=emp.get('family_status', 'EE'),
            costs=emp.get('costs', {}),
            winner=emp.get('winner', ''),
            insight=emp.get('insight', ''),
            metal_plan_details=emp.get('metal_plan_details', {}),
            member_breakdowns=emp.get('member_breakdowns', {}),
            current_total_monthly=emp.get('current_total_monthly', 0),
            renewal_total_monthly=emp.get('renewal_total_monthly', 0),
            use_ee_rate_only=emp.get('use_ee_rate_only', False),
            has_renewal_data=has_renewal_data,
            contribution_strategy=emp.get('contribution_strategy', ''),
        ))

    generator = EmployeeExamplesSlideGenerator(
        client_name=client_name,
        plan_config=plan_config,
        include_qr_links=include_qr_links
    )
    return generator.generate(employees)


if __name__ == "__main__":
    # Test with sample data
    sample_employees = [
        {
            'label': 'Youngest Employee',
            'name': 'John Smith',
            'age': 25,
            'tier': 'Employee Only',
            'location': 'Philadelphia, PA (Rating Area 7)',
            'family_ages': [],
            'costs': {
                'Current': {'employee': 150, 'employer': 450},
                'Renewal': {'employee': 175, 'employer': 525},
                'ICHRA Bronze': {'employee': 50, 'employer': 350},
                'ICHRA Silver': {'employee': 100, 'employer': 400},
                'ICHRA Gold': {'employee': 150, 'employer': 450},
                'Cooperative': {'employee': 80, 'employer': 320},
            },
            'winner': 'ICHRA Bronze',
            'insight': 'This young employee saves $75/month with ICHRA Bronze compared to renewal.',
            'metal_plan_details': {
                'Bronze': {'plan_name': 'Keystone HMO Bronze', 'premium': 400, 'deductible': 7500, 'moop': 9200},
                'Silver': {'plan_name': 'Keystone HMO Silver', 'premium': 500, 'deductible': 4500, 'moop': 9200},
                'Gold': {'plan_name': 'Keystone HMO Gold', 'premium': 600, 'deductible': 1500, 'moop': 6500},
            },
            'renewal_total_monthly': 700,
        },
    ]

    pptx_buffer = generate_employee_examples_pptx(sample_employees, "Test Company")

    with open("/tmp/employee_examples_test.pptx", "wb") as f:
        f.write(pptx_buffer.read())

    print("✓ Test PPTX generated: /tmp/employee_examples_test.pptx")
