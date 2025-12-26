"""
PPTX Template Filler

Fills a PowerPoint template with data from ProposalData by replacing
{placeholder} text with actual values.

Workflow:
1. Design your template in PowerPoint with {placeholder} syntax
2. This code loads the template and replaces placeholders with values
3. Static slides (no placeholders) pass through unchanged

This approach separates design (PowerPoint) from data (Python).
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from pathlib import Path
import re
from typing import Dict, Optional

from pptx_generator import ProposalData


class PPTXTemplateFiller:
    """Fill a PowerPoint template with data from ProposalData"""

    TEMPLATE_PATH = Path(__file__).parent / 'templates' / 'glove_proposal_template.pptx'

    def __init__(self, data: ProposalData, template_path: Optional[str] = None):
        """
        Initialize filler with proposal data.

        Args:
            data: ProposalData instance with all proposal fields
            template_path: Optional custom template path (defaults to TEMPLATE_PATH)
        """
        self.data = data
        self.template_path = Path(template_path) if template_path else self.TEMPLATE_PATH
        self.placeholder_pattern = re.compile(r'\{(\w+)\}')

    def _get_replacements(self) -> Dict[str, str]:
        """
        Build replacement dictionary from ProposalData.

        Returns:
            Dict mapping placeholder names to formatted values
        """
        d = self.data

        # Calculate derived values
        savings_pct = 0
        if d.total_renewal_cost > 0:
            savings_pct = (d.annual_savings_vs_renewal / d.total_renewal_cost) * 100

        return {
            # === Cover Slide ===
            'client_name': d.client_name,
            'client_name_upper': d.client_name.upper(),
            'renewal_percentage': f"{d.renewal_percentage:.0f}",
            'total_renewal_cost': f"${d.total_renewal_cost:,.0f}",
            'total_renewal_cost_m': f"${d.total_renewal_cost/1_000_000:.1f}M" if d.total_renewal_cost >= 1_000_000 else f"${d.total_renewal_cost/1_000:.0f}K",

            # === Market/Census Data ===
            'employee_count': str(d.employee_count),
            'avg_monthly_premium': f"${d.avg_monthly_premium:,.0f}",
            'covered_lives': str(d.covered_lives),
            'total_employees': str(d.total_employees),
            'total_dependents': str(d.total_dependents),
            'total_states': str(d.total_states),
            'per_life_monthly': f"${d.per_life_monthly:,.0f}",

            # === Fit Score ===
            'fit_score': str(d.fit_score),

            # Category scores (for radar chart labels if needed)
            'score_cost_advantage': str(d.category_scores.get('cost_advantage', 0)),
            'score_market_readiness': str(d.category_scores.get('market_readiness', 0)),
            'score_workforce_fit': str(d.category_scores.get('workforce_fit', 0)),
            'score_geographic': str(d.category_scores.get('geographic_complexity', 0)),
            'score_employee_exp': str(d.category_scores.get('employee_experience', 0)),
            'score_admin': str(d.category_scores.get('admin_readiness', 0)),

            # === Demographics ===
            'avg_employee_age': f"{d.avg_employee_age:.0f}",
            'age_range': f"{d.age_range_min}-{d.age_range_max}",
            'age_range_min': str(d.age_range_min),
            'age_range_max': str(d.age_range_max),

            # === Family Status Breakdown ===
            'family_ee': str(d.family_status_breakdown.get('EE', 0)),
            'family_es': str(d.family_status_breakdown.get('ES', 0)),
            'family_ec': str(d.family_status_breakdown.get('EC', 0)),
            'family_f': str(d.family_status_breakdown.get('F', 0)),

            # === Current Costs ===
            'current_total_monthly': f"${d.current_total_monthly:,.0f}",
            'current_total_annual': f"${d.current_total_annual:,.0f}",
            'current_er_monthly': f"${d.current_er_monthly:,.0f}",
            'current_er_annual': f"${d.current_er_annual:,.0f}",
            'current_ee_monthly': f"${d.current_ee_monthly:,.0f}",
            'current_ee_annual': f"${d.current_ee_annual:,.0f}",

            # === Proposed ICHRA Costs ===
            'proposed_er_monthly': f"${d.proposed_er_monthly:,.0f}",
            'proposed_er_annual': f"${d.proposed_er_annual:,.0f}",

            # === Workflow Slide Values ===
            'renewal_monthly': f"${d.renewal_monthly:,.0f}",
            'ichra_monthly': f"${d.ichra_monthly:,.0f}",
            'current_to_renewal_diff': f"+${d.current_to_renewal_diff_monthly:,.0f}",
            'current_to_renewal_pct': f"+{d.current_to_renewal_pct:.1f}%",
            'renewal_to_ichra_diff': f"-${d.renewal_to_ichra_diff_monthly:,.0f}",
            'renewal_to_ichra_pct': f"-{d.renewal_to_ichra_pct:.1f}%",

            # === Savings ===
            'annual_savings': f"${d.annual_savings:,.0f}",
            'annual_savings_vs_renewal': f"${d.annual_savings_vs_renewal:,.0f}",
            'savings_percentage': f"{d.savings_percentage:.0f}%",
            'savings_pct_calculated': f"{savings_pct:.0f}%",

            # === Healthcare Burden ===
            'healthcare_burden': f"${d.additional_healthcare_burden:,.0f}",
            'healthcare_burden_k': f"${d.additional_healthcare_burden/1000:,.0f}K",

            # === ICHRA Analysis Allowance Levels ===
            # $450 level
            'allowance_450': '$450',
            'allowance_450_annual': f"${450 * d.employee_count * 12:,.0f}",
            'allowance_450_savings': f"${d.total_renewal_cost - (450 * d.employee_count * 12):,.0f}" if d.total_renewal_cost > 0 else "$0",

            # $600 level
            'allowance_600': '$600',
            'allowance_600_annual': f"${600 * d.employee_count * 12:,.0f}",
            'allowance_600_savings': f"${d.total_renewal_cost - (600 * d.employee_count * 12):,.0f}" if d.total_renewal_cost > 0 else "$0",

            # $750 level
            'allowance_750': '$750',
            'allowance_750_annual': f"${750 * d.employee_count * 12:,.0f}",
            'allowance_750_savings': f"${d.total_renewal_cost - (750 * d.employee_count * 12):,.0f}" if d.total_renewal_cost > 0 else "$0",

            # === Top States (for geographic slide) ===
            'state_1': d.top_states[0]['state'] if len(d.top_states) > 0 else '',
            'state_1_count': str(d.top_states[0]['count']) if len(d.top_states) > 0 else '',
            'state_2': d.top_states[1]['state'] if len(d.top_states) > 1 else '',
            'state_2_count': str(d.top_states[1]['count']) if len(d.top_states) > 1 else '',
            'state_3': d.top_states[2]['state'] if len(d.top_states) > 2 else '',
            'state_3_count': str(d.top_states[2]['count']) if len(d.top_states) > 2 else '',
            'state_4': d.top_states[3]['state'] if len(d.top_states) > 3 else '',
            'state_4_count': str(d.top_states[3]['count']) if len(d.top_states) > 3 else '',
            'state_5': d.top_states[4]['state'] if len(d.top_states) > 4 else '',
            'state_5_count': str(d.top_states[4]['count']) if len(d.top_states) > 4 else '',
        }

    def _replace_text_in_shape(self, shape, replacements: Dict[str, str]) -> None:
        """
        Replace {placeholders} in a shape's text while preserving formatting.

        Args:
            shape: PowerPoint shape object
            replacements: Dict mapping placeholder names to values
        """
        if not shape.has_text_frame:
            return

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text
                new_text = original_text

                # Find and replace all {placeholder} patterns
                for match in self.placeholder_pattern.finditer(original_text):
                    placeholder = match.group(1)  # e.g., "client_name"
                    if placeholder in replacements:
                        new_text = new_text.replace(
                            f"{{{placeholder}}}",
                            replacements[placeholder]
                        )

                if new_text != original_text:
                    run.text = new_text

    def _process_shape(self, shape, replacements: Dict[str, str]) -> None:
        """
        Process a shape and its children for placeholder replacement.

        Args:
            shape: PowerPoint shape object
            replacements: Dict mapping placeholder names to values
        """
        # Handle text in the shape itself
        self._replace_text_in_shape(shape, replacements)

        # Handle grouped shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for subshape in shape.shapes:
                self._process_shape(subshape, replacements)

        # Handle tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    self._replace_text_in_shape(cell, replacements)

    def _process_slide(self, slide, replacements: Dict[str, str]) -> None:
        """
        Process all shapes in a slide for placeholder replacement.

        Args:
            slide: PowerPoint slide object
            replacements: Dict mapping placeholder names to values
        """
        for shape in slide.shapes:
            self._process_shape(shape, replacements)

    def generate(self) -> BytesIO:
        """
        Generate filled PPTX from template.

        Returns:
            BytesIO buffer containing the filled PPTX

        Raises:
            FileNotFoundError: If template file doesn't exist
        """
        if not self.template_path.exists():
            raise FileNotFoundError(
                f"Template not found: {self.template_path}\n"
                f"Please create the template with {placeholder} syntax."
            )

        # Load template
        prs = Presentation(str(self.template_path))

        # Build replacement dict
        replacements = self._get_replacements()

        # Process each slide
        for slide in prs.slides:
            self._process_slide(slide, replacements)

        # Save to buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    def list_placeholders(self) -> Dict[str, str]:
        """
        Get all available placeholders and their current values.
        Useful for template design reference.

        Returns:
            Dict of placeholder names to formatted values
        """
        return self._get_replacements()


def generate_pptx_proposal(data: ProposalData, template_path: Optional[str] = None) -> BytesIO:
    """
    Convenience function to generate PPTX proposal.

    Args:
        data: ProposalData instance
        template_path: Optional custom template path

    Returns:
        BytesIO buffer containing the filled PPTX
    """
    filler = PPTXTemplateFiller(data, template_path)
    return filler.generate()


if __name__ == "__main__":
    # Test/demo with sample data
    from pptx_generator import ProposalData

    test_data = ProposalData(
        client_name="Ecentria",
        renewal_percentage=40,
        total_renewal_cost=1_300_000,
        employee_count=226,
        avg_monthly_premium=559,
        covered_lives=377,
        total_employees=226,
        total_dependents=151,
        fit_score=87,
        category_scores={
            'cost_advantage': 85,
            'market_readiness': 80,
            'workforce_fit': 90,
            'geographic_complexity': 85,
            'employee_experience': 75,
            'admin_readiness': 90,
        },
        total_states=18,
        top_states=[
            {'state': 'IL', 'count': 180},
            {'state': 'AZ', 'count': 11},
            {'state': 'UT', 'count': 5},
            {'state': 'TX', 'count': 4},
            {'state': 'IN', 'count': 4},
        ],
        family_status_breakdown={'EE': 155, 'ES': 26, 'EC': 14, 'F': 31},
        current_total_monthly=203048,
        current_total_annual=2436581,
        current_er_monthly=148219,
        current_er_annual=1778631,
        current_ee_monthly=54829,
        current_ee_annual=657950,
        per_life_monthly=538,
        proposed_er_monthly=115935,
        proposed_er_annual=1391216,
        renewal_monthly=84585,
        ichra_monthly=47752,
        current_to_renewal_diff_monthly=6806,
        current_to_renewal_pct=8.75,
        renewal_to_ichra_diff_monthly=36833,
        renewal_to_ichra_pct=43.5,
        annual_savings=387415,
        annual_savings_vs_renewal=441988,
        savings_percentage=21.8,
        additional_healthcare_burden=280000,
        avg_employee_age=40.3,
        age_range_min=21,
        age_range_max=73,
    )

    print("Available placeholders for template design:")
    print("=" * 50)

    filler = PPTXTemplateFiller(test_data)
    placeholders = filler.list_placeholders()

    for name, value in sorted(placeholders.items()):
        print(f"  {{{name}}} → {value}")

    print("\n" + "=" * 50)
    print("To use: Create glove_proposal_template.pptx with these placeholders")
    print("Place in: templates/glove_proposal_template.pptx")
