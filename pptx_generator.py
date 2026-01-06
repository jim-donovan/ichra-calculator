"""
PowerPoint Proposal Generator for ICHRA Calculator

Uses python-pptx to populate a GLOVE template with census data,
financial analysis, and Fit Score calculations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import Dict, List, Optional, Any
from io import BytesIO
from pathlib import Path
from dataclasses import dataclass, field
import pandas as pd
import re
import copy

from financial_calculator import FinancialSummaryCalculator

# Template paths
TEMPLATE_DIR = Path(__file__).parent / 'templates'
WORKFLOW_SLIDE_PATH = TEMPLATE_DIR / 'ICHRA_evaluation_workflow.pptx'


@dataclass
class ProposalData:
    """Data container for all proposal fields"""

    # Slide 1: Cover
    client_name: str = "ABC Company"
    renewal_percentage: float = 0.0
    total_renewal_cost: float = 0.0

    # Slide 2: Market Overview
    employee_count: int = 0
    avg_monthly_premium: float = 0.0

    # Slide 3: Fit Score
    fit_score: int = 0
    category_scores: Dict[str, int] = field(default_factory=dict)

    # Slide 5: Cost burden
    additional_healthcare_burden: float = 0.0
    total_annual_salaries: float = 0.0

    # Slide 8: Geographic distribution
    total_states: int = 0
    top_states: List[Dict[str, Any]] = field(default_factory=list)

    # Slide 9: Census demographics
    covered_lives: int = 0
    total_employees: int = 0
    total_dependents: int = 0
    total_spouses: int = 0
    total_children: int = 0
    avg_employee_age: float = 0.0
    avg_dependent_age: float = 0.0
    age_range_min: int = 0
    age_range_max: int = 0
    family_status_breakdown: Dict[str, int] = field(default_factory=dict)
    current_er_monthly: float = 0.0
    current_ee_monthly: float = 0.0
    current_total_monthly: float = 0.0
    current_er_annual: float = 0.0
    current_ee_annual: float = 0.0
    current_total_annual: float = 0.0
    per_life_monthly: float = 0.0

    # Slide 10: ICHRA analysis
    proposed_er_monthly: float = 0.0
    proposed_ee_monthly: float = 0.0
    proposed_er_annual: float = 0.0
    proposed_ee_annual: float = 0.0
    annual_savings: float = 0.0
    savings_percentage: float = 0.0

    # Projected Renewal ER calculations (applying current ER/EE split to 2026 renewal)
    er_contribution_pct: float = 0.0  # Current ER% of total (e.g., 0.596 = 59.6%)
    ee_contribution_pct: float = 0.0  # Current EE% of total
    projected_er_monthly_2026: float = 0.0  # renewal_monthly * er_pct
    projected_er_annual_2026: float = 0.0
    projected_ee_monthly_2026: float = 0.0  # renewal_monthly * ee_pct
    projected_ee_annual_2026: float = 0.0

    # Standardized comparisons
    savings_vs_renewal_er: float = 0.0  # projected_er_annual_2026 - proposed_er_annual (positive = savings)
    savings_vs_renewal_er_pct: float = 0.0
    delta_vs_current_er: float = 0.0  # proposed_er_annual - current_er_annual (positive = costs more)
    delta_vs_current_er_pct: float = 0.0

    # ICHRA Evaluation Workflow slide
    renewal_monthly: float = 0.0
    ichra_monthly: float = 0.0
    current_to_renewal_diff_monthly: float = 0.0
    current_to_renewal_pct: float = 0.0
    renewal_to_ichra_diff_monthly: float = 0.0
    renewal_to_ichra_pct: float = 0.0
    annual_increase: float = 0.0  # Renewal annual - Current annual
    annual_savings_vs_renewal: float = 0.0  # Renewal annual - ICHRA annual

    # Contribution level analysis
    contribution_450_annual: float = 0.0
    contribution_450_savings: float = 0.0
    contribution_600_annual: float = 0.0
    contribution_600_savings: float = 0.0
    contribution_750_annual: float = 0.0
    contribution_750_savings: float = 0.0

    # Slide 11: Plans by state
    plans_by_state: Dict[str, int] = field(default_factory=dict)

    @classmethod
    def from_session_state(cls, session_state) -> 'ProposalData':
        """
        Build ProposalData from Streamlit session_state.

        Args:
            session_state: st.session_state object

        Returns:
            Populated ProposalData instance
        """
        data = cls()

        census_df = session_state.get('census_df')
        dependents_df = session_state.get('dependents_df')
        financial_summary = session_state.get('financial_summary', {})
        contribution_settings = session_state.get('contribution_settings', {})

        if census_df is None or census_df.empty:
            return data

        # Basic counts
        data.total_employees = len(census_df)
        data.employee_count = data.total_employees

        # Calculate covered lives: employees + dependents (matches Census Input page)
        if dependents_df is not None and not dependents_df.empty:
            data.total_dependents = len(dependents_df)
        else:
            data.total_dependents = 0

        data.covered_lives = data.total_employees + data.total_dependents

        # Count spouses vs children from dependents_df
        if dependents_df is not None and not dependents_df.empty:
            if 'relationship' in dependents_df.columns:
                data.total_spouses = len(dependents_df[dependents_df['relationship'] == 'spouse'])
                data.total_children = len(dependents_df[dependents_df['relationship'] == 'child'])
            else:
                # Estimate from family status in census
                if 'family_status' in census_df.columns:
                    family_status = census_df['family_status'].fillna('EE')
                    data.total_spouses = family_status.isin(['ES', 'F']).sum()
                    data.total_children = max(0, data.total_dependents - data.total_spouses)

        # Age statistics
        age_col = 'age' if 'age' in census_df.columns else 'employee_age'
        if age_col in census_df.columns:
            ages = census_df[age_col].dropna()
            if not ages.empty:
                data.avg_employee_age = round(ages.mean(), 1)
                data.age_range_min = int(ages.min())
                data.age_range_max = int(ages.max())

        if dependents_df is not None and 'age' in dependents_df.columns:
            dep_ages = dependents_df['age'].dropna()
            if not dep_ages.empty:
                data.avg_dependent_age = round(dep_ages.mean(), 1)

        # Geographic distribution
        if 'state' in census_df.columns:
            state_counts = census_df['state'].value_counts()
            data.total_states = len(state_counts)
            data.top_states = [
                {'state': state, 'count': int(count)}
                for state, count in state_counts.head(5).items()
            ]

        # Family status breakdown
        if 'family_status' in census_df.columns:
            data.family_status_breakdown = census_df['family_status'].value_counts().to_dict()

        # Current costs from census
        if 'current_er_monthly' in census_df.columns:
            er_sum = census_df['current_er_monthly'].sum()
            if pd.notna(er_sum):
                data.current_er_monthly = float(er_sum)
                data.current_er_annual = data.current_er_monthly * 12

        if 'current_ee_monthly' in census_df.columns:
            ee_sum = census_df['current_ee_monthly'].sum()
            if pd.notna(ee_sum):
                data.current_ee_monthly = float(ee_sum)
                data.current_ee_annual = data.current_ee_monthly * 12

        data.current_total_monthly = data.current_er_monthly + data.current_ee_monthly
        data.current_total_annual = data.current_total_monthly * 12

        if data.covered_lives > 0:
            data.per_life_monthly = data.current_total_monthly / data.covered_lives

        # ICHRA projections - check strategy_results first (from Contribution Evaluation page)
        strategy_results = session_state.get('strategy_results', {})
        strategy_result = strategy_results.get('result') or strategy_results.get('current') or {}

        if strategy_result:
            data.proposed_er_monthly = strategy_result.get('total_monthly', 0)
            data.proposed_er_annual = strategy_result.get('total_annual', 0)
        else:
            # Fallback to financial_summary (from LCSP Analysis page)
            results = financial_summary.get('results', {})
            if results:
                data.proposed_er_monthly = results.get('total_monthly', 0)
                data.proposed_er_annual = results.get('total_annual', 0)

        # Calculate ER/EE contribution split percentages (aggregate - kept for backwards compatibility)
        if data.current_total_monthly > 0:
            data.er_contribution_pct = data.current_er_monthly / data.current_total_monthly
            data.ee_contribution_pct = data.current_ee_monthly / data.current_total_monthly
        else:
            # Default to typical 60/40 ER/EE split if no current data
            data.er_contribution_pct = 0.60
            data.ee_contribution_pct = 0.40

        # Detect per-tier contribution pattern (percentage vs flat-rate)
        contribution_pattern = None
        if census_df is not None and not census_df.empty:
            from utils import ContributionComparison
            if ContributionComparison.has_individual_contributions(census_df):
                contribution_pattern = ContributionComparison.detect_contribution_pattern(census_df)

        # Calculate savings vs current ER (legacy field for backwards compatibility)
        if data.proposed_er_annual > 0 and data.current_er_annual > 0:
            data.annual_savings = data.current_er_annual - data.proposed_er_annual
            data.savings_percentage = (data.annual_savings / data.current_er_annual) * 100

            # Delta vs current ER (positive = ICHRA costs more than current)
            data.delta_vs_current_er = data.proposed_er_annual - data.current_er_annual
            data.delta_vs_current_er_pct = (data.delta_vs_current_er / data.current_er_annual) * 100

        # Average monthly premium per employee (from current group plan)
        if data.total_employees > 0 and data.current_total_monthly > 0:
            data.avg_monthly_premium = data.current_total_monthly / data.total_employees

        # 2026 Renewal data from financial_summary or calculate from census
        renewal_monthly = financial_summary.get('renewal_monthly', 0) or 0

        # If no renewal data in financial_summary, try to calculate from census 2026 Premium column
        if renewal_monthly == 0 and census_df is not None:
            from financial_calculator import FinancialSummaryCalculator
            projected_data = FinancialSummaryCalculator.calculate_projected_2026_total(census_df)
            if projected_data.get('has_data'):
                renewal_monthly = projected_data.get('total_monthly', 0)

        # No fallback - require actual renewal data
        # If no renewal data available, leave as 0 (will show N/A in proposals)

        data.renewal_monthly = renewal_monthly
        data.ichra_monthly = data.proposed_er_monthly

        if renewal_monthly > 0:
            data.total_renewal_cost = renewal_monthly * 12  # Annual

            # PROJECT the 2026 ER/EE contribution using detected tier-level pattern
            # This applies the detected contribution pattern (% or flat-rate) per tier
            if contribution_pattern is not None and census_df is not None:
                # Apply pattern to each employee and sum
                census_with_projections = ContributionComparison.apply_contribution_pattern(
                    census_df, contribution_pattern
                )
                data.projected_er_monthly_2026 = census_with_projections['projected_2026_er'].sum()
                data.projected_ee_monthly_2026 = census_with_projections['projected_2026_ee'].sum()
            else:
                # Fallback: use aggregate ER/EE split ratio
                data.projected_er_monthly_2026 = renewal_monthly * data.er_contribution_pct
                data.projected_ee_monthly_2026 = renewal_monthly * data.ee_contribution_pct

            data.projected_er_annual_2026 = data.projected_er_monthly_2026 * 12
            data.projected_ee_annual_2026 = data.projected_ee_monthly_2026 * 12

            # PRIMARY SALES COMPARISON: ICHRA vs Projected Renewal ER
            # This is the comparison that matters - "what you'd pay at renewal vs what ICHRA costs"
            if data.proposed_er_annual > 0 and data.projected_er_annual_2026 > 0:
                data.savings_vs_renewal_er = data.projected_er_annual_2026 - data.proposed_er_annual
                data.savings_vs_renewal_er_pct = (data.savings_vs_renewal_er / data.projected_er_annual_2026) * 100

            # Calculate renewal increase % vs current total (2026 Premium is total renewal)
            if data.current_total_monthly > 0:
                increase = renewal_monthly - data.current_total_monthly
                data.renewal_percentage = (increase / data.current_total_monthly) * 100

                # ICHRA Evaluation Workflow calculations
                data.current_to_renewal_diff_monthly = increase
                data.current_to_renewal_pct = data.renewal_percentage
                data.annual_increase = increase * 12

            # Renewal to ICHRA comparison (vs total renewal, for context)
            if data.proposed_er_monthly > 0:
                ichra_diff = renewal_monthly - data.proposed_er_monthly
                data.renewal_to_ichra_diff_monthly = ichra_diff
                data.renewal_to_ichra_pct = (ichra_diff / renewal_monthly) * 100 if renewal_monthly > 0 else 0
                data.annual_savings_vs_renewal = ichra_diff * 12

        # Calculate total annual salaries and healthcare burden (30% of salaries)
        if 'monthly_income' in census_df.columns:
            monthly_salaries = census_df['monthly_income'].fillna(0)
            # Parse currency values if needed
            if monthly_salaries.dtype == 'object':
                from utils import parse_currency
                monthly_salaries = monthly_salaries.apply(lambda x: parse_currency(str(x)) if pd.notna(x) else 0)
            data.total_annual_salaries = float(monthly_salaries.sum()) * 12
            data.additional_healthcare_burden = data.total_annual_salaries * 0.30  # 30% of salaries

        return data

    def validate(self) -> tuple:
        """
        Validate comparison data for internal consistency.

        Returns:
            (errors: List[str], warnings: List[str])
        """
        errors = []
        warnings = []

        # Check cost totals sum correctly
        calculated_total = self.current_ee_annual + self.current_er_annual
        if abs(calculated_total - self.current_total_annual) > 1:
            errors.append(
                f"Current costs don't sum: EE({self.current_ee_annual:,.0f}) + ER({self.current_er_annual:,.0f}) "
                f"!= Total({self.current_total_annual:,.0f})"
            )

        # Check ER percentage is reasonable (typically 40-90%)
        if self.er_contribution_pct > 0:
            if self.er_contribution_pct < 0.4:
                warnings.append(f"Unusual ER contribution percentage: {self.er_contribution_pct*100:.1f}% (< 40%)")
            elif self.er_contribution_pct > 0.9:
                warnings.append(f"Unusual ER contribution percentage: {self.er_contribution_pct*100:.1f}% (> 90%)")

        # Check projected renewal ER calculation
        if self.renewal_monthly > 0 and self.er_contribution_pct > 0:
            expected_projected_er = self.renewal_monthly * self.er_contribution_pct * 12
            if abs(expected_projected_er - self.projected_er_annual_2026) > 1:
                errors.append(
                    f"Projected renewal ER mismatch: expected {expected_projected_er:,.0f}, "
                    f"got {self.projected_er_annual_2026:,.0f}"
                )

        # Check ICHRA vs renewal makes sense
        if self.savings_vs_renewal_er < 0 and self.proposed_er_annual > 0:
            warnings.append(
                f"ICHRA costs MORE than projected renewal ER by ${abs(self.savings_vs_renewal_er):,.0f} - "
                "verify contribution strategy"
            )

        # Check renewal is higher than current (typical scenario)
        if self.renewal_monthly > 0 and self.current_total_monthly > 0:
            if self.renewal_monthly < self.current_total_monthly:
                warnings.append(
                    f"Renewal ({self.renewal_monthly:,.0f}/mo) is LOWER than current "
                    f"({self.current_total_monthly:,.0f}/mo) - unusual"
                )

        # Check employee count
        if self.employee_count <= 0:
            errors.append("Employee count is zero or negative")

        # Check for zero proposed ICHRA when strategy should be applied
        if self.proposed_er_annual == 0 and self.employee_count > 0:
            warnings.append("Proposed ICHRA is $0 - has a contribution strategy been applied?")

        return errors, warnings


class ProposalGenerator:
    """Generate GLOVE ICHRA proposals from PowerPoint template"""

    # Template placeholder patterns for each slide
    # Maps slide index (0-based) to {placeholder_pattern: data_field}
    SLIDE_MAPPINGS = {
        0: {  # Slide 1: Cover
            'ECENTRIA': 'client_name',
            '40': 'renewal_percentage_display',
            '$1.3': 'total_cost_display',
        },
        1: {  # Slide 2: Market Overview
            '45': 'employee_count',
            '$559': 'avg_monthly_premium_display',
        },
        2: {  # Slide 3: Fit Score
            '87': 'fit_score',
        },
        4: {  # Slide 5: Cost burden
            '226': 'employee_count',
            '$280K': 'healthcare_burden_display',
        },
        7: {  # Slide 8: Geographic
            '266': 'employee_count',
            '18': 'total_states',
        },
        8: {  # Slide 9: Census demographics
            '377': 'covered_lives',
            '226': 'total_employees',
            '151': 'total_dependents',
        },
        9: {  # Slide 10: ICHRA Analysis
            '$387,414.78': 'annual_savings_display',
        },
    }

    def __init__(self, template_path: str):
        """
        Initialize with template PowerPoint file.

        Args:
            template_path: Path to GLOVE template .pptx file
        """
        self.template_path = Path(template_path)
        self.prs: Optional[Presentation] = None

    def load_template(self) -> None:
        """Load the PowerPoint template file"""
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {self.template_path}")
        self.prs = Presentation(str(self.template_path))

    def discover_shapes(self, slide_index: int) -> List[Dict]:
        """
        Discover all shapes on a slide for mapping purposes.

        Args:
            slide_index: 0-based slide index

        Returns:
            List of shape info dicts
        """
        if self.prs is None:
            raise RuntimeError("Template not loaded. Call load_template() first.")

        if slide_index >= len(self.prs.slides):
            return []

        slide = self.prs.slides[slide_index]
        shapes_info = []

        for shape in slide.shapes:
            info = {
                'name': shape.name,
                'shape_type': str(shape.shape_type),
                'has_text_frame': shape.has_text_frame,
                'text': '',
                'left': shape.left,
                'top': shape.top,
            }

            if shape.has_text_frame:
                info['text'] = shape.text_frame.text[:100]

            shapes_info.append(info)

        return shapes_info

    def find_shape_by_text(self, slide, text_pattern: str):
        """
        Find a shape by searching for text content pattern.

        Args:
            slide: PowerPoint slide object
            text_pattern: Text to search for (exact or partial match)

        Returns:
            Shape object if found, None otherwise
        """
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_text = shape.text_frame.text
                if text_pattern in shape_text:
                    return shape
        return None

    def find_shapes_containing_text(self, slide, text_pattern: str) -> List:
        """
        Find all shapes containing specific text.

        Args:
            slide: PowerPoint slide object
            text_pattern: Text to search for

        Returns:
            List of matching shape objects
        """
        matches = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                if text_pattern in shape.text_frame.text:
                    matches.append(shape)
        return matches

    def replace_text_in_shape(self, shape, old_text: str, new_text: str) -> bool:
        """
        Replace text in a shape while attempting to preserve formatting.

        Args:
            shape: PowerPoint shape object
            old_text: Text to find
            new_text: Replacement text

        Returns:
            True if replacement was made, False otherwise
        """
        if not shape.has_text_frame:
            return False

        replaced = False

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, str(new_text))
                    replaced = True

        # Fallback: replace in full paragraph text if run-level didn't work
        if not replaced:
            for paragraph in shape.text_frame.paragraphs:
                full_text = ''.join(run.text for run in paragraph.runs)
                if old_text in full_text:
                    # Clear runs and set new text on first run
                    new_full_text = full_text.replace(old_text, str(new_text))
                    if paragraph.runs:
                        # Keep first run's formatting, clear others
                        first_run = paragraph.runs[0]
                        first_run.text = new_full_text
                        for run in paragraph.runs[1:]:
                            run.text = ''
                        replaced = True

        return replaced

    def populate_slide_1(self, slide, data: ProposalData) -> None:
        """Populate Slide 1: Cover/Title slide"""
        # Replace client name
        self.replace_text_in_shape(
            self.find_shape_by_text(slide, 'ECENTRIA') or self.find_shape_by_text(slide, 'is at a'),
            'ECENTRIA',
            data.client_name.upper()
        )

        # Replace renewal percentage (the "40" in "40%")
        shape = self.find_shape_by_text(slide, '40')
        if shape:
            self.replace_text_in_shape(shape, '40', str(int(data.renewal_percentage)))

        # Replace total cost
        total_cost_m = data.total_renewal_cost / 1_000_000
        shape = self.find_shape_by_text(slide, '$1.3')
        if shape:
            self.replace_text_in_shape(shape, '$1.3', f'${total_cost_m:.1f}')

    def populate_slide_2(self, slide, data: ProposalData) -> None:
        """Populate Slide 2: Market Overview"""
        # Employee count
        shape = self.find_shape_by_text(slide, '45')
        if shape and '45' in shape.text_frame.text:
            self.replace_text_in_shape(shape, '45', str(data.employee_count))

        # Avg monthly premium
        shape = self.find_shape_by_text(slide, '$559')
        if shape:
            self.replace_text_in_shape(shape, '$559', f'${data.avg_monthly_premium:,.0f}')

    def populate_slide_3(self, slide, data: ProposalData) -> None:
        """Populate Slide 3: Fit Score"""
        # Main fit score (87)
        shape = self.find_shape_by_text(slide, '87')
        if shape:
            self.replace_text_in_shape(shape, '87', str(data.fit_score))

    def populate_slide_5(self, slide, data: ProposalData) -> None:
        """Populate Slide 5: Cost burden"""
        # Employee count
        shape = self.find_shape_by_text(slide, '226')
        if shape:
            self.replace_text_in_shape(shape, '226', str(data.employee_count))

        # Healthcare burden
        burden_k = data.additional_healthcare_burden / 1000
        shape = self.find_shape_by_text(slide, '$280K')
        if shape:
            self.replace_text_in_shape(shape, '$280K', f'${burden_k:,.0f}K')

    def populate_slide_8(self, slide, data: ProposalData) -> None:
        """Populate Slide 8: Geographic distribution"""
        # Total employees
        shape = self.find_shape_by_text(slide, '266')
        if shape:
            self.replace_text_in_shape(shape, '266', str(data.employee_count))

        # Total states
        shape = self.find_shape_by_text(slide, '18')
        if shape and 'STATES' in slide.shapes[slide.shapes.index(shape) + 1].text_frame.text if hasattr(slide.shapes[slide.shapes.index(shape) + 1], 'text_frame') else True:
            self.replace_text_in_shape(shape, '18', str(data.total_states))

        # Top 5 states - find and replace state names and counts
        state_name_map = {
            'illinois': 0, 'ILLINOIS': 0,
            'ARIZONA': 1,
            'UTAH': 2,
            'TEXAS': 3,
            'INDIANA': 4,
        }

        # Try to find and update each state placeholder
        for placeholder_state, idx in state_name_map.items():
            if idx < len(data.top_states):
                state_data = data.top_states[idx]
                shape = self.find_shape_by_text(slide, placeholder_state)
                if shape:
                    self.replace_text_in_shape(shape, placeholder_state, state_data['state'].upper())

        # Update state employee counts
        count_placeholders = ['180', '11', '5', '4', '4']
        for idx, placeholder in enumerate(count_placeholders):
            if idx < len(data.top_states):
                shape = self.find_shape_by_text(slide, placeholder)
                if shape:
                    self.replace_text_in_shape(shape, placeholder, str(data.top_states[idx]['count']))

    def populate_slide_9(self, slide, data: ProposalData) -> None:
        """Populate Slide 9: Census demographics"""
        replacements = [
            ('377', str(data.covered_lives)),
            ('226', str(data.total_employees)),
            ('151', str(data.total_dependents)),
            ('155', str(data.family_status_breakdown.get('EE', 0))),
            ('26', str(data.family_status_breakdown.get('ES', 0))),
            ('14', str(data.family_status_breakdown.get('EC', 0))),
            ('31', str(data.family_status_breakdown.get('F', 0))),
            ('40.3', f'{data.avg_employee_age:.1f}'),
            ('21-73', f'{data.age_range_min}-{data.age_range_max}'),
            ('$2,436,581', f'${data.current_total_annual:,.0f}'),
            ('$203,048/mo', f'${data.current_total_monthly:,.0f}/mo'),
            ('$538/mo', f'${data.per_life_monthly:,.0f}/mo'),
            ('$1,778,631', f'${data.current_er_annual:,.0f}'),
            ('$657,950', f'${data.current_ee_annual:,.0f}'),
        ]

        for old_text, new_text in replacements:
            shapes = self.find_shapes_containing_text(slide, old_text)
            for shape in shapes:
                self.replace_text_in_shape(shape, old_text, new_text)

    def populate_slide_10(self, slide, data: ProposalData) -> None:
        """Populate Slide 10: ICHRA Analysis

        Key update: Now shows savings vs Projected Renewal ER (not vs current ER)
        This is the PRIMARY sales comparison - avoiding the renewal hit.
        """
        savings_display = abs(data.savings_vs_renewal_er) if data.savings_vs_renewal_er != 0 else abs(data.annual_savings)

        replacements = [
            # Primary savings metric - vs Renewal ER
            ('$387,414.78', f'${savings_display:,.2f}'),
            # Current ER costs
            ('$1,778,631.00', f'${data.current_er_annual:,.2f}'),
            ('$148,219.25', f'${data.current_er_monthly:,.2f}'),
            # Current EE costs
            ('$657,949.44', f'${data.current_ee_annual:,.2f}'),
            ('$54,829.12', f'${data.current_ee_monthly:,.2f}'),
            # Total ER Cost
            ('$1,391,216.22', f'${data.proposed_er_annual:,.2f}'),
            ('$115,934.68', f'${data.proposed_er_monthly:,.2f}'),
        ]

        for old_text, new_text in replacements:
            shapes = self.find_shapes_containing_text(slide, old_text)
            for shape in shapes:
                self.replace_text_in_shape(shape, old_text, new_text)

    def populate_workflow_slide(self, slide, data: ProposalData) -> None:
        """Populate ICHRA Evaluation Workflow slide (appended at end)"""
        # Main value replacements
        replacements = [
            # Monthly values
            ('$77,779', f'${data.current_total_monthly:,.0f}'),
            ('$84,585', f'${data.renewal_monthly:,.0f}'),
            ('$47,752', f'${data.ichra_monthly:,.0f}'),
            # Annual values
            ('$1,015,018', f'${data.total_renewal_cost:,.0f}'),
            ('$573,030', f'${data.proposed_er_annual:,.0f}'),
            # Savings
            ('$441,988', f'${data.annual_savings_vs_renewal:,.0f}'),
        ]

        for old_text, new_text in replacements:
            shapes = self.find_shapes_containing_text(slide, old_text)
            for shape in shapes:
                self.replace_text_in_shape(shape, old_text, new_text)

        # Delta badges (inside groups)
        delta_replacements = [
            ('+$6,806', f'+${data.current_to_renewal_diff_monthly:,.0f}'),
            ('-$36,833', f'-${data.renewal_to_ichra_diff_monthly:,.0f}'),
            ('-44%', f'-{data.renewal_to_ichra_pct:.0f}%'),
        ]

        # Search in groups for delta badges
        for shape in slide.shapes:
            if shape.shape_type == 6:  # GROUP
                self._replace_in_group(shape, delta_replacements)

    def _replace_in_group(self, group, replacements: List[tuple]) -> None:
        """Recursively replace text in grouped shapes"""
        for shape in group.shapes:
            if shape.has_text_frame:
                for old_text, new_text in replacements:
                    if old_text in shape.text_frame.text:
                        self.replace_text_in_shape(shape, old_text, new_text)
            if shape.shape_type == 6:  # Nested group
                self._replace_in_group(shape, replacements)

    def append_workflow_slide(self) -> None:
        """Append the ICHRA Evaluation Workflow slide from separate template.

        Uses proper slide duplication that preserves images, shapes, and relationships
        by duplicating the slide XML and copying all related media parts.
        """
        if not WORKFLOW_SLIDE_PATH.exists():
            raise FileNotFoundError(f"Workflow slide not found: {WORKFLOW_SLIDE_PATH}")

        # Load the workflow slide presentation (pre-resized to match glove template)
        workflow_prs = Presentation(str(WORKFLOW_SLIDE_PATH))

        if len(workflow_prs.slides) == 0:
            return None

        source_slide = workflow_prs.slides[0]

        # Duplicate slide using internal XML copy with relationship handling
        new_slide = self._duplicate_slide(source_slide, workflow_prs)

        return new_slide

    def _duplicate_slide(self, source_slide, source_prs):
        """Duplicate a slide from source presentation to destination, preserving all content.

        Args:
            source_slide: Source slide to copy
            source_prs: Source presentation (needed for media access)

        Returns:
            The new slide added to self.prs
        """
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.parts.image import ImagePart
        from pptx.opc.packuri import PackURI
        from lxml import etree

        # Calculate scale factor if source and destination have different dimensions
        scale_factor = self.prs.slide_width / source_prs.slide_width

        # Add blank slide using a layout
        try:
            blank_layout = self.prs.slide_layouts[6]  # Usually blank
        except IndexError:
            blank_layout = self.prs.slide_layouts[0]

        new_slide = self.prs.slides.add_slide(blank_layout)

        # Copy the slide's shape tree (spTree) content
        # First, clear the new slide's shapes (except placeholders we want to keep)
        sp_tree = new_slide.shapes._spTree
        # Remove all existing shapes from new slide
        for sp in list(sp_tree):
            if sp.tag.endswith('}sp') or sp.tag.endswith('}pic') or sp.tag.endswith('}grpSp'):
                sp_tree.remove(sp)

        # Build a mapping of old relationship IDs to new ones
        rId_map = {}

        # Track used image indices to avoid conflicts
        used_indices = set()
        for part in self.prs.part.package.iter_parts():
            partname = str(part.partname)
            if '/ppt/media/image' in partname:
                # Extract index from partname like /ppt/media/image5.png
                match = re.search(r'image(\d+)', partname)
                if match:
                    used_indices.add(int(match.group(1)))

        next_idx = max(used_indices, default=0) + 1

        # Copy media relationships (images) from source to destination
        for rel in source_slide.part.rels.values():
            if rel.reltype == RT.IMAGE:
                try:
                    # Get the image part from source
                    image_part = rel.target_part
                    image_blob = image_part.blob
                    content_type = image_part.content_type
                    ext = image_part.partname.ext  # e.g., '.png', '.svg'

                    # Ensure ext has a dot prefix
                    if not ext.startswith('.'):
                        ext = '.' + ext

                    # Create a unique partname
                    new_partname = f'/ppt/media/image{next_idx}{ext}'
                    next_idx += 1

                    # Create a new ImagePart with the image data using ImagePart.load
                    # Signature: load(partname, content_type, package, blob)
                    new_part = ImagePart.load(
                        PackURI(new_partname),
                        content_type,
                        self.prs.part.package,
                        image_blob
                    )

                    # Create relationship from slide to image
                    new_rId = new_slide.part.relate_to(new_part, RT.IMAGE)
                    rId_map[rel.rId] = new_rId
                except Exception:
                    # Skip images that fail to copy
                    pass

        # Copy shapes from source slide
        source_sp_tree = source_slide.shapes._spTree
        for child in source_sp_tree:
            tag_name = etree.QName(child.tag).localname
            if tag_name in ('sp', 'pic', 'grpSp', 'cxnSp', 'graphicFrame'):
                # Deep copy the element
                new_elem = copy.deepcopy(child)

                # Scale positions, sizes, and text properties if needed
                if scale_factor != 1.0:
                    self._scale_element(new_elem, scale_factor)

                # Update relationship IDs in the copied element
                self._update_rIds(new_elem, rId_map)

                # Insert into destination
                sp_tree.append(new_elem)

        return new_slide

    def _scale_element(self, element, scale_factor: float):
        """Scale all dimensional values in an XML element.

        Handles positions, sizes, font sizes, line spacing, margins, etc.

        Args:
            element: XML element to scale
            scale_factor: Scale factor to apply (e.g., 0.5 for half size)
        """
        from lxml import etree
        a_ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        # Attributes that contain EMU values (positions, sizes)
        emu_attrs = ['x', 'y', 'cx', 'cy', 'w', 'l', 'r', 't', 'b',
                     'lIns', 'rIns', 'tIns', 'bIns',  # Text insets
                     'indent', 'marL', 'marR']  # Margins

        # Attributes that contain centipoints (font sizes, spacing)
        centipoint_attrs = ['sz', 'val']  # sz for font size, val for spacing points

        for elem in element.iter():
            # Scale EMU attributes
            for attr in emu_attrs:
                if attr in elem.attrib:
                    try:
                        old_val = int(elem.attrib[attr])
                        elem.attrib[attr] = str(int(old_val * scale_factor))
                    except ValueError:
                        pass

            # Scale font sizes (sz attribute on rPr, defRPr, endParaRPr elements)
            tag_local = etree.QName(elem.tag).localname
            if tag_local in ('rPr', 'defRPr', 'endParaRPr', 'lvl1pPr', 'lvl2pPr'):
                if 'sz' in elem.attrib:
                    try:
                        old_sz = int(elem.attrib['sz'])
                        elem.attrib['sz'] = str(int(old_sz * scale_factor))
                    except ValueError:
                        pass

            # Scale line spacing (spcPts element)
            if tag_local == 'spcPts':
                if 'val' in elem.attrib:
                    try:
                        old_val = int(elem.attrib['val'])
                        elem.attrib['val'] = str(int(old_val * scale_factor))
                    except ValueError:
                        pass

            # Scale line widths
            if tag_local == 'ln':
                if 'w' in elem.attrib:
                    try:
                        old_w = int(elem.attrib['w'])
                        elem.attrib['w'] = str(int(old_w * scale_factor))
                    except ValueError:
                        pass

    def _update_rIds(self, element, rId_map: Dict[str, str]):
        """Update relationship IDs in an XML element based on mapping.

        Args:
            element: XML element to update
            rId_map: Dict mapping old rIds to new rIds
        """
        # Namespaces used in PowerPoint XML
        r_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

        # Update r:embed and r:link attributes
        for attr_name in ['embed', 'link']:
            full_attr = f'{r_ns}{attr_name}'
            for elem in element.iter():
                if full_attr in elem.attrib:
                    old_rId = elem.attrib[full_attr]
                    if old_rId in rId_map:
                        elem.attrib[full_attr] = rId_map[old_rId]

    def populate_all_slides(self, proposal_data: ProposalData, include_workflow: bool = True) -> None:
        """
        Populate all slides with proposal data.

        Args:
            proposal_data: ProposalData object with all calculated values
            include_workflow: Whether to append and populate the ICHRA Evaluation Workflow slide
        """
        if self.prs is None:
            raise RuntimeError("Template not loaded. Call load_template() first.")

        slides = self.prs.slides

        # Populate each slide
        if len(slides) >= 1:
            self.populate_slide_1(slides[0], proposal_data)

        if len(slides) >= 2:
            self.populate_slide_2(slides[1], proposal_data)

        if len(slides) >= 3:
            self.populate_slide_3(slides[2], proposal_data)

        if len(slides) >= 5:
            self.populate_slide_5(slides[4], proposal_data)

        if len(slides) >= 8:
            self.populate_slide_8(slides[7], proposal_data)

        if len(slides) >= 9:
            self.populate_slide_9(slides[8], proposal_data)

        if len(slides) >= 10:
            self.populate_slide_10(slides[9], proposal_data)

        # Append and populate ICHRA Evaluation Workflow slide at the end
        if include_workflow and WORKFLOW_SLIDE_PATH.exists():
            workflow_slide = self.append_workflow_slide()
            if workflow_slide:
                self.populate_workflow_slide(workflow_slide, proposal_data)

    def generate(self) -> BytesIO:
        """
        Generate final PowerPoint and return as BytesIO buffer.

        Returns:
            BytesIO buffer containing the .pptx file
        """
        if self.prs is None:
            raise RuntimeError("Template not loaded. Call load_template() first.")

        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer

    def save(self, output_path: str) -> None:
        """
        Save PowerPoint to file.

        Args:
            output_path: Path to save the file
        """
        if self.prs is None:
            raise RuntimeError("Template not loaded. Call load_template() first.")

        self.prs.save(output_path)


def generate_proposal(
    template_path: str,
    proposal_data: ProposalData
) -> BytesIO:
    """
    Convenience function to generate a proposal.

    Args:
        template_path: Path to template file
        proposal_data: ProposalData with values to populate

    Returns:
        BytesIO buffer with generated PowerPoint
    """
    generator = ProposalGenerator(template_path)
    generator.load_template()
    generator.populate_all_slides(proposal_data)
    return generator.generate()


if __name__ == "__main__":
    # Test with sample data
    from pathlib import Path

    template_path = Path(__file__).parent / 'templates' / 'glove_template.pptx'

    if template_path.exists():
        # Create sample data
        data = ProposalData(
            client_name="Acme Corp",
            renewal_percentage=35.0,
            total_renewal_cost=2_500_000,
            employee_count=150,
            avg_monthly_premium=485,
            fit_score=82,
            covered_lives=275,
            total_employees=150,
            total_dependents=125,
            total_states=5,
            top_states=[
                {'state': 'CA', 'count': 80},
                {'state': 'TX', 'count': 35},
                {'state': 'NY', 'count': 20},
                {'state': 'FL', 'count': 10},
                {'state': 'IL', 'count': 5},
            ],
            current_er_annual=1_500_000,
            proposed_er_annual=1_100_000,
            annual_savings=400_000,
        )

        # Generate
        buffer = generate_proposal(str(template_path), data)
        print(f"Generated PowerPoint: {len(buffer.getvalue())} bytes")
        print("Test complete!")
    else:
        print(f"Template not found at: {template_path}")
