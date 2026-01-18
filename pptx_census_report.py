"""
Census Analysis Report Slide Generator

Generates a single PowerPoint slide with census demographics
matching the Figma design for Census Analysis Report.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, Optional
from io import BytesIO
from pathlib import Path
from dataclasses import dataclass
from datetime import datetime
import pandas as pd


# Color scheme matching the Figma design
COLORS = {
    # Primary brand colors
    'teal_dark': RGBColor(0x11, 0x5E, 0x59),      # #115E59 - Header title
    'teal_border': RGBColor(0x0F, 0x76, 0x6E),    # #0F766E - Header border
    'teal_accent': RGBColor(0x5E, 0xD4, 0xC4),    # #5ED4C4 - Section divider

    # Text colors
    'text_primary': RGBColor(0x10, 0x18, 0x28),   # #101828 - Main text (matches other templates)
    'text_secondary': RGBColor(0x6B, 0x72, 0x80), # #6B7280 - Labels
    'text_muted': RGBColor(0x9C, 0xA3, 0xAF),     # #9CA3AF - Footer
    'text_body': RGBColor(0x4B, 0x55, 0x63),      # #4B5563 - Body text
    'text_label': RGBColor(0x37, 0x41, 0x51),     # #374151 - Labels in cards
    'section_header': RGBColor(0x36, 0x41, 0x51), # #364151 - Section headers and card values

    # Card colors
    'card_bg': RGBColor(0xF9, 0xFA, 0xFB),        # #F9FAFB - Card background

    # Dependent Overview cards - brand palette (teal, golden, burnt sienna, brown)
    # All pass WCAG AA with text_body (#4B5563) labels: 6.9-7.3:1 contrast
    'dep_card_1_bg': RGBColor(0xE6, 0xFF, 0xFA),     # #E6FFFA - teal tint
    'dep_card_1_border': RGBColor(0x37, 0xBE, 0xAE), # #37BEAE - Glove brand teal
    'dep_card_2_bg': RGBColor(0xF5, 0xF3, 0xFF),     # #F5F3FF - violet-50
    'dep_card_2_border': RGBColor(0x7C, 0x3A, 0xED), # #7C3AED - violet-600
    'dep_card_3_bg': RGBColor(0xEF, 0xF6, 0xFF),     # #EFF6FF - blue-50
    'dep_card_3_border': RGBColor(0x1D, 0x4E, 0xD8), # #1D4ED8 - blue-700
    'dep_card_4_bg': RGBColor(0xEE, 0xF2, 0xFF),     # #EEF2FF - indigo-50
    'dep_card_4_border': RGBColor(0x4F, 0x46, 0xE5), # #4F46E5 - indigo-600

    # Progress bar colors
    'progress_bg': RGBColor(0xE5, 0xE7, 0xEB),    # #E5E7EB - Progress bar bg
    'progress_cyan': RGBColor(0x38, 0xBF, 0xB1),  # #38BFB1 - Children progress
    'progress_brown': RGBColor(0x78, 0x35, 0x0F), # #78350F - Spouses progress

    # Base
    'white': RGBColor(0xFF, 0xFF, 0xFF),
}

# Slide dimensions (standard widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Decorative image paths
DECORATIVES_DIR = Path(__file__).parent / "decoratives"
CORNER_IMAGE = DECORATIVES_DIR / "glove-tile-corner.png"
BANNER_IMAGE = DECORATIVES_DIR / "PPT_header.png"


@dataclass
class CensusReportData:
    """Data container for Census Analysis Report slide"""

    # Top metrics
    total_employees: int = 0
    total_dependents: int = 0
    covered_lives: int = 0
    states_count: int = 0

    # Age statistics (all covered lives)
    avg_age_all: float = 0.0
    median_age_all: float = 0.0
    min_age_all: int = 0
    max_age_all: int = 0

    # Employee demographics
    emp_avg_age: float = 0.0
    emp_median_age: float = 0.0
    emp_min_age: int = 0
    emp_max_age: int = 0

    # Dependent statistics
    dep_avg_age: float = 0.0
    dep_median_age: float = 0.0
    dep_min_age: int = 0
    dep_max_age: int = 0
    dep_children_count: int = 0
    dep_spouses_count: int = 0

    # Coverage burden
    coverage_burden: float = 0.0

    # Client info
    client_name: str = ""
    generated_date: str = ""

    @classmethod
    def from_census_data(cls, employees_df: pd.DataFrame,
                         dependents_df: pd.DataFrame = None,
                         client_name: str = "") -> 'CensusReportData':
        """
        Build CensusReportData from census DataFrames.

        Args:
            employees_df: Employee census DataFrame
            dependents_df: Dependents DataFrame (optional)
            client_name: Client/company name

        Returns:
            Populated CensusReportData instance
        """
        data = cls()
        data.client_name = client_name
        data.generated_date = datetime.now().strftime("%m.%d.%y")

        if employees_df is None or employees_df.empty:
            return data

        # Employee counts and ages
        data.total_employees = len(employees_df)

        if 'age' in employees_df.columns:
            emp_ages = employees_df['age'].dropna()
            if not emp_ages.empty:
                data.emp_avg_age = round(emp_ages.mean(), 1)
                data.emp_median_age = round(emp_ages.median(), 1)
                data.emp_min_age = int(emp_ages.min())
                data.emp_max_age = int(emp_ages.max())

        # State count
        state_col = 'state' if 'state' in employees_df.columns else 'Home State'
        if state_col in employees_df.columns:
            data.states_count = employees_df[state_col].nunique()

        # Dependent statistics
        all_ages = list(employees_df['age'].dropna()) if 'age' in employees_df.columns else []

        if dependents_df is not None and not dependents_df.empty:
            data.total_dependents = len(dependents_df)

            if 'age' in dependents_df.columns:
                dep_ages = dependents_df['age'].dropna()
                if not dep_ages.empty:
                    data.dep_avg_age = round(dep_ages.mean(), 1)
                    data.dep_median_age = round(dep_ages.median(), 1)
                    data.dep_min_age = int(dep_ages.min())
                    data.dep_max_age = int(dep_ages.max())
                    all_ages.extend(list(dep_ages))

            # Count by relationship type
            if 'relationship' in dependents_df.columns:
                rel_counts = dependents_df['relationship'].value_counts()
                data.dep_children_count = rel_counts.get('Child', 0) + rel_counts.get('child', 0)
                data.dep_spouses_count = rel_counts.get('Spouse', 0) + rel_counts.get('spouse', 0)
            else:
                # Estimate based on typical patterns if relationship not specified
                data.dep_children_count = data.total_dependents
                data.dep_spouses_count = 0

        # Covered lives and coverage burden
        data.covered_lives = data.total_employees + data.total_dependents
        if data.total_employees > 0:
            data.coverage_burden = round(data.covered_lives / data.total_employees, 2)

        # All covered lives age statistics
        if all_ages:
            import numpy as np
            data.avg_age_all = round(np.mean(all_ages), 1)
            data.median_age_all = round(np.median(all_ages), 1)
            data.min_age_all = int(min(all_ages))
            data.max_age_all = int(max(all_ages))

        return data


class CensusReportSlideGenerator:
    """Generate Census Analysis Report PowerPoint slide"""

    def __init__(self):
        """Initialize with a new blank presentation"""
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    def _add_text_box(self, slide, left, top, width, height, text: str,
                      font_size: int = 14, font_name: str = 'Poppins',
                      color: RGBColor = None, bold: bool = False,
                      align: PP_ALIGN = PP_ALIGN.LEFT):
        """Add a text box with formatting"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        p.alignment = align
        return txBox

    def _add_metric_card(self, slide, left, top, width, label: str, value: str,
                         unit: str = None):
        """Add a metric card (top row style)"""
        # Label
        self._add_text_box(
            slide, left, top, width, Inches(0.3),
            label, font_size=14, color=COLORS['text_secondary'],
            bold=True, align=PP_ALIGN.CENTER
        )

        # Value centered (no unit display - "yrs" is obvious)
        self._add_text_box(
            slide, left, top + Inches(0.35), width, Inches(0.6),
            value, font_size=32, color=COLORS['text_primary'],
            align=PP_ALIGN.CENTER
        )

    def _add_colored_card(self, slide, left, top, width, height,
                          label: str, value: str, bg_color: RGBColor,
                          border_color: RGBColor):
        """Add a colored card with left border accent"""
        # Background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = bg_color

        # Left border accent (thin rectangle)
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height
        )
        border.fill.solid()
        border.fill.fore_color.rgb = border_color
        border.line.fill.background()

        # Label (font size 10) - use text_body for WCAG AA accessibility
        self._add_text_box(
            slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.3), Inches(0.25),
            label, font_size=10, color=COLORS['text_body'], bold=True
        )

        # Value (font size 24)
        self._add_text_box(
            slide, left + Inches(0.2), top + Inches(0.4), width - Inches(0.3), Inches(0.4),
            value, font_size=24, color=COLORS['section_header'], bold=True
        )

    def _add_progress_bar(self, slide, left, top, width, height,
                          fill_pct: float, fill_color: RGBColor):
        """Add a progress bar"""
        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['progress_bg']
        bg.line.fill.background()

        # Fill
        if fill_pct > 0:
            fill_width = width * (fill_pct / 100)
            fill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, fill_width, height
            )
            fill.fill.solid()
            fill.fill.fore_color.rgb = fill_color
            fill.line.fill.background()

    def _add_section_card(self, slide, left, top, width, height):
        """Add a section card background"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['card_bg']
        shape.line.fill.background()
        return shape

    def _add_decorative_images(self, slide):
        """Add decorative images - banner at top, corner tile bottom-right"""
        # Banner at top (same as Cooperative Health slide)
        if BANNER_IMAGE.exists():
            try:
                slide.shapes.add_picture(
                    str(BANNER_IMAGE),
                    left=Inches(0),
                    top=Inches(0),
                    width=SLIDE_WIDTH,
                    height=Inches(0.25)
                )
            except Exception:
                pass  # Skip if image can't be added

        # Bottom-right corner decoration - exact PPT position from user spec
        if CORNER_IMAGE.exists():
            try:
                slide.shapes.add_picture(
                    str(CORNER_IMAGE),
                    Inches(12.58),  # Horizontal position from top-left
                    Inches(6.78),   # Vertical position from top-left
                    width=Inches(0.75),
                    height=Inches(0.75)
                )
            except Exception:
                pass  # Skip if image can't be added

    def create_slide(self, data: CensusReportData):
        """Create the Census Analysis Report slide"""
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)

        # Add decorative images (if available)
        self._add_decorative_images(slide)

        # Margins
        margin_left = Inches(0.5)
        margin_top = Inches(0.4)
        content_width = SLIDE_WIDTH - Inches(1.0)

        # === HEADER ===
        # Title
        self._add_text_box(
            slide, margin_left, margin_top, Inches(6), Inches(0.6),
            "Census analysis", font_size=32,
            color=COLORS['teal_dark'], bold=True
        )

        # Client name (right aligned, font size 24, wider box for longer names)
        client_display = data.client_name if data.client_name else "Client"
        self._add_text_box(
            slide, SLIDE_WIDTH - Inches(5.5), margin_top + Inches(0.05), Inches(5), Inches(0.5),
            client_display, font_size=24, color=COLORS['text_secondary'],
            bold=True, align=PP_ALIGN.RIGHT
        )

        # Header border line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, margin_left, margin_top + Inches(0.65),
            content_width, Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS['teal_border']
        line.line.fill.background()

        # === TOP METRICS ROW ===
        metrics_top = margin_top + Inches(0.9)
        metric_width = Inches(2.8)
        metric_spacing = Inches(0.3)

        metrics = [
            ("TOTAL EMPLOYEES", str(data.total_employees)),
            ("TOTAL DEPENDENTS", str(data.total_dependents)),
            ("COVERED LIVES", str(data.covered_lives)),
            ("STATES", str(data.states_count)),
        ]

        for i, (label, value) in enumerate(metrics):
            x = margin_left + i * (metric_width + metric_spacing)
            self._add_metric_card(slide, x, metrics_top, metric_width, label, value)

        # === COVERAGE STATS ROW ===
        stats_top = metrics_top + Inches(1.1)
        stat_width = Inches(3.8)

        age_stats = [
            ("AVG AGE OF COVERED LIVES", f"{data.avg_age_all:.1f}", "yrs"),
            ("MEDIAN AGE OF COVERED LIVES", f"{data.median_age_all:.1f}", "yrs"),
            ("AGE RANGE OF COVERED LIVES", f"{data.min_age_all} - {data.max_age_all}", "yrs"),
        ]

        for i, (label, value, unit) in enumerate(age_stats):
            x = margin_left + i * (stat_width + Inches(0.2))
            self._add_metric_card(slide, x, stats_top, stat_width, label, value, unit)

        # Section divider line
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, margin_left, stats_top + Inches(1.15),
            content_width, Inches(0.04)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = COLORS['teal_accent']
        divider.line.fill.background()

        # === MAIN CONTENT GRID ===
        main_top = stats_top + Inches(1.4)
        left_col_width = Inches(4.0)
        right_col_left = margin_left + left_col_width + Inches(0.5)
        right_col_width = content_width - left_col_width - Inches(0.5)

        # --- LEFT COLUMN: Employee Demographics ---
        self._add_text_box(
            slide, margin_left, main_top, left_col_width, Inches(0.4),
            "Employee demographics", font_size=20,
            color=COLORS['section_header'], bold=True
        )

        demo_top = main_top + Inches(0.5)
        line_height = Inches(0.35)

        emp_items = [
            ("Total employees", str(data.total_employees), True),
            ("Average age", f"{data.emp_avg_age:.1f} years", False),
            ("Median age", f"{data.emp_median_age:.1f} years", False),
        ]

        for i, (label, value, bold) in enumerate(emp_items):
            y = demo_top + i * line_height
            # Wider label box for "Total employees" to prevent wrap
            self._add_text_box(
                slide, margin_left, y, Inches(2.0), line_height,
                label, font_size=14, color=COLORS['text_body']
            )
            self._add_text_box(
                slide, margin_left + Inches(2.0), y, Inches(1.5), line_height,
                value, font_size=14, color=COLORS['text_primary'], bold=bold
            )

        # Age range subsection
        range_top = demo_top + 3 * line_height + Inches(0.15)
        self._add_text_box(
            slide, margin_left, range_top, left_col_width, Inches(0.3),
            "Age range", font_size=16, color=COLORS['section_header'], bold=True
        )

        range_items = [
            ("Youngest", f"{data.emp_min_age} years"),
            ("Oldest", f"{data.emp_max_age} years"),
        ]

        for i, (label, value) in enumerate(range_items):
            y = range_top + Inches(0.35) + i * line_height
            self._add_text_box(
                slide, margin_left + Inches(0.2), y, Inches(1.5), line_height,
                label, font_size=14, color=COLORS['text_body']
            )
            self._add_text_box(
                slide, margin_left + Inches(1.7), y, Inches(1.5), line_height,
                value, font_size=14, color=COLORS['text_primary']
            )

        # --- RIGHT COLUMN: Dependent Overview ---
        self._add_text_box(
            slide, right_col_left, main_top, right_col_width, Inches(0.4),
            "Dependent overview", font_size=20,
            color=COLORS['section_header'], bold=True
        )

        # 4 Colored Cards
        cards_top = main_top + Inches(0.5)
        card_width = Inches(1.9)
        card_height = Inches(0.85)
        card_spacing = Inches(0.15)

        # Calculate dependent percentages for display
        total_deps = data.total_dependents if data.total_dependents > 0 else 1
        children_pct = (data.dep_children_count / total_deps) * 100 if total_deps > 0 else 0
        spouses_pct = (data.dep_spouses_count / total_deps) * 100 if total_deps > 0 else 0

        dep_cards = [
            ("TOTAL DEPENDENTS", str(data.total_dependents), COLORS['dep_card_1_bg'], COLORS['dep_card_1_border']),
            ("COVERAGE BURDEN", f"{data.coverage_burden:.2f}:1", COLORS['dep_card_2_bg'], COLORS['dep_card_2_border']),
            ("AVERAGE AGE", f"{data.dep_avg_age:.1f}", COLORS['dep_card_3_bg'], COLORS['dep_card_3_border']),
            ("AGE RANGE", f"{data.dep_min_age}-{data.dep_max_age}" if data.dep_max_age > 0 else "--", COLORS['dep_card_4_bg'], COLORS['dep_card_4_border']),
        ]

        for i, (label, value, bg, border) in enumerate(dep_cards):
            x = right_col_left + i * (card_width + card_spacing)
            self._add_colored_card(slide, x, cards_top, card_width, card_height,
                                   label, value, bg, border)

        # Bottom sections
        sections_top = cards_top + card_height + Inches(0.25)
        section_width = (right_col_width - Inches(0.3)) / 2
        section_height = Inches(1.6)

        # By Relationship section
        self._add_section_card(slide, right_col_left, sections_top,
                               section_width, section_height)

        self._add_text_box(
            slide, right_col_left + Inches(0.2), sections_top + Inches(0.15),
            section_width - Inches(0.4), Inches(0.3),
            "By relationship", font_size=16, color=COLORS['section_header'], bold=True
        )

        # Children row
        child_y = sections_top + Inches(0.5)
        self._add_text_box(
            slide, right_col_left + Inches(0.2), child_y, Inches(1.2), Inches(0.25),
            "Children", font_size=14, color=COLORS['text_label']
        )
        self._add_text_box(
            slide, right_col_left + section_width - Inches(1.4), child_y,
            Inches(1.2), Inches(0.25),
            f"{data.dep_children_count} ({children_pct:.0f}%)",
            font_size=14, color=COLORS['text_body'], align=PP_ALIGN.RIGHT
        )
        self._add_progress_bar(
            slide, right_col_left + Inches(0.2), child_y + Inches(0.3),
            section_width - Inches(0.4), Inches(0.15),
            children_pct, COLORS['progress_cyan']
        )

        # Spouses row
        spouse_y = sections_top + Inches(1.0)
        self._add_text_box(
            slide, right_col_left + Inches(0.2), spouse_y, Inches(1.2), Inches(0.25),
            "Spouses", font_size=14, color=COLORS['text_label']
        )
        self._add_text_box(
            slide, right_col_left + section_width - Inches(1.4), spouse_y,
            Inches(1.2), Inches(0.25),
            f"{data.dep_spouses_count} ({spouses_pct:.0f}%)",
            font_size=14, color=COLORS['text_body'], align=PP_ALIGN.RIGHT
        )
        self._add_progress_bar(
            slide, right_col_left + Inches(0.2), spouse_y + Inches(0.3),
            section_width - Inches(0.4), Inches(0.15),
            spouses_pct, COLORS['progress_brown']
        )

        # Dependents Age Statistics section
        stats_left = right_col_left + section_width + Inches(0.3)
        self._add_section_card(slide, stats_left, sections_top,
                               section_width, section_height)

        self._add_text_box(
            slide, stats_left + Inches(0.2), sections_top + Inches(0.15),
            section_width - Inches(0.4), Inches(0.3),
            "Dependents age statistics", font_size=16,
            color=COLORS['section_header'], bold=True
        )

        age_stat_items = [
            ("Average", f"{data.dep_avg_age:.1f} years"),
            ("Median", f"{data.dep_median_age:.1f} years"),
        ]

        for i, (label, value) in enumerate(age_stat_items):
            y = sections_top + Inches(0.55) + i * Inches(0.4)
            self._add_text_box(
                slide, stats_left + Inches(0.2), y, Inches(1.2), Inches(0.3),
                label, font_size=16, color=COLORS['text_body']
            )
            self._add_text_box(
                slide, stats_left + section_width - Inches(1.8), y,
                Inches(1.6), Inches(0.3),
                value, font_size=16, color=COLORS['text_primary'],
                bold=True, align=PP_ALIGN.RIGHT
            )

        # === FOOTER ===
        # Coverage burden explanation (font size 10) - positioned at bottom
        footer_top = SLIDE_HEIGHT - Inches(0.4)
        burden_text = (
            f"Coverage burden = (employees + dependents) / employees = "
            f"({data.total_employees} + {data.total_dependents}) / {data.total_employees} = "
            f"{data.coverage_burden:.2f} covered lives per employee"
        )
        self._add_text_box(
            slide, margin_left, footer_top, content_width, Inches(0.25),
            burden_text, font_size=10, color=COLORS['text_muted']
        )

        return slide

    def generate(self) -> BytesIO:
        """Generate the PowerPoint file and return as BytesIO"""
        output = BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output


def generate_census_report_slide(data: CensusReportData) -> BytesIO:
    """
    Convenience function to generate a Census Report slide.

    Args:
        data: CensusReportData with values to populate

    Returns:
        BytesIO buffer with generated PowerPoint
    """
    generator = CensusReportSlideGenerator()
    generator.create_slide(data)
    return generator.generate()


def generate_census_report_from_session(session_state) -> BytesIO:
    """
    Generate Census Report slide from Streamlit session state.

    Args:
        session_state: st.session_state object

    Returns:
        BytesIO buffer with generated PowerPoint
    """
    employees_df = session_state.get('census_df')
    dependents_df = session_state.get('dependents_df')
    client_name = session_state.get('client_name', '')

    data = CensusReportData.from_census_data(
        employees_df=employees_df,
        dependents_df=dependents_df,
        client_name=client_name
    )

    return generate_census_report_slide(data)


def get_census_report_filename(client_name: str = "") -> str:
    """
    Generate filename with client name and timestamp.

    Args:
        client_name: Optional client name

    Returns:
        Filename string like "ClientName_census_report_20260117_143052.pptx"
    """
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    if client_name:
        # Clean client name for filename
        clean_name = "".join(c for c in client_name if c.isalnum() or c in (' ', '-', '_')).strip()
        clean_name = clean_name.replace(' ', '_')
        return f"{clean_name}_census_report_{timestamp}.pptx"
    return f"census_report_{timestamp}.pptx"
