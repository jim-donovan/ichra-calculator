"""
Plan Comparison Slide Generator

Generates a single PowerPoint slide with the Plan Comparison table
that matches the design from Page 9.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from typing import List, Optional
from io import BytesIO
from pathlib import Path
from dataclasses import dataclass, field


# Color scheme matching the comparison table design
# Using more saturated colors (-100/-200 variants) for better visibility in PPT
COLORS = {
    # Current Plan - Gray
    'current_bg': RGBColor(0xF3, 0xF4, 0xF6),      # #F3F4F6 - gray-100
    'current_text': RGBColor(0x37, 0x41, 0x51),    # #374151
    'current_header': RGBColor(0x37, 0x41, 0x51),  # #374151

    # Comparison colors - more saturated for PPT visibility
    'better_bg': RGBColor(0xD1, 0xFA, 0xE5),       # #D1FAE5 - emerald-100
    'better_text': RGBColor(0x05, 0x96, 0x69),     # #059669
    'similar_bg': RGBColor(0xE8, 0xF1, 0xFD),      # #E8F1FD - cobalt-50 (light blue)
    'similar_text': RGBColor(0x1D, 0x4E, 0xD8),    # #1D4ED8 - blue-700
    'worse_bg': RGBColor(0xFE, 0xE2, 0xE2),        # #FEE2E2 - red-100
    'worse_text': RGBColor(0xDC, 0x26, 0x26),      # #DC2626

    # Metal level colors
    'bronze': RGBColor(0xD9, 0x77, 0x06),          # #D97706
    'silver': RGBColor(0x63, 0x66, 0xF1),          # #6366F1
    'gold': RGBColor(0x05, 0x96, 0x69),            # #059669
    'platinum': RGBColor(0x7C, 0x3A, 0xED),        # #7C3AED

    # Neutral
    'row_label_bg': RGBColor(0xF9, 0xFA, 0xFB),    # #F9FAFB
    'row_label_text': RGBColor(0x37, 0x41, 0x51),  # #374151
    'section_bg': RGBColor(0xF9, 0xFA, 0xFB),      # #F9FAFB
    'section_text': RGBColor(0x6B, 0x72, 0x80),    # #6B7280
    'detail_text': RGBColor(0x6B, 0x72, 0x80),     # #6B7280
    'white': RGBColor(0xFF, 0xFF, 0xFF),           # #FFFFFF
    'black': RGBColor(0x11, 0x18, 0x27),           # #111827
    'border': RGBColor(0xE5, 0xE7, 0xEB),          # #E5E7EB
}

# Slide dimensions (standard widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Decorative image paths
DECORATIVE_IMAGES_DIR = Path(__file__).parent.parent / "glove-design" / "✍️ Design + Crit 01"
CORNER_IMAGE = DECORATIVE_IMAGES_DIR / "glove-tile-corner.png"
EDGE_IMAGE = DECORATIVE_IMAGES_DIR / "glove-tile-edge-h-fade.png"


@dataclass
class PlanColumnData:
    """Data for a single plan column"""
    plan_name: str
    issuer_name: str = ""
    plan_type: tuple = ("", 'similar')  # (value, comparison) - same as UI
    metal_level: str = ""
    hsa_eligible: tuple = (False, 'similar')  # (value, comparison)
    actuarial_value: Optional[float] = None
    is_current: bool = False

    # Premium data
    age_21_premium: Optional[float] = None
    total_premium: Optional[float] = None
    avg_premium: Optional[float] = None
    affordable_contribution: Optional[float] = None
    contribution_range: str = ""

    # Premium comparisons (calculated same way as UI)
    age_21_premium_comparison: Optional[str] = None
    total_premium_comparison: Optional[str] = None

    # Current plan specific - for showing current/renewal pairs
    current_age_21_premium: Optional[float] = None
    renewal_age_21_premium: Optional[float] = None
    current_total_premium: Optional[float] = None
    renewal_total_premium: Optional[float] = None
    current_avg_premium: Optional[float] = None
    renewal_avg_premium: Optional[float] = None

    # Benefits (value, comparison result)
    individual_deductible: tuple = (0.0, 'similar')
    family_deductible: tuple = (None, 'similar')
    individual_oop_max: tuple = (0.0, 'similar')
    family_oop_max: tuple = (None, 'similar')
    coinsurance_pct: tuple = (20, 'similar')


@dataclass
class PlanComparisonSlideData:
    """Data container for Plan Comparison slide"""
    plans: List[PlanColumnData] = field(default_factory=list)
    employee_count: int = 0
    avg_age: float = 0.0
    footnote: str = ""
    client_name: str = ""


def get_metal_color(metal_level: str) -> RGBColor:
    """Get color for metal level"""
    level_lower = metal_level.lower()
    if 'bronze' in level_lower:
        return COLORS['bronze']
    elif 'silver' in level_lower:
        return COLORS['silver']
    elif 'gold' in level_lower:
        return COLORS['gold']
    elif 'platinum' in level_lower:
        return COLORS['platinum']
    return COLORS['current_header']


class PlanComparisonSlideGenerator:
    """Generate Plan Comparison PowerPoint slide"""

    def __init__(self):
        """Initialize with a new blank presentation"""
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    def _set_cell_fill(self, cell, color: RGBColor):
        """Set cell background color"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    def _set_cell_text(self, cell, text: str, color: RGBColor,
                       font_size: int = 11, bold: bool = False,
                       align: PP_ALIGN = PP_ALIGN.CENTER):
        """Set cell text with formatting"""
        cell.text = str(text) if text is not None else "—"
        cell.text_frame.paragraphs[0].alignment = align
        cell.text_frame.paragraphs[0].font.name = 'Poppins'
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.text_frame.paragraphs[0].font.color.rgb = color
        cell.text_frame.paragraphs[0].font.bold = bold
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _add_cell_line(self, cell, text: str, color: RGBColor,
                       font_size: int = 9, bold: bool = False,
                       align: PP_ALIGN = PP_ALIGN.CENTER):
        """Add additional line to cell"""
        p = cell.text_frame.add_paragraph()
        p.text = str(text) if text is not None else ""
        p.alignment = align
        p.font.name = 'Poppins'
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold

    def _format_currency(self, value) -> str:
        """Format value as currency"""
        if value is None:
            return "—"
        try:
            return f"${float(value):,.0f}"
        except (ValueError, TypeError):
            return "—"

    def create_slide(self, data: PlanComparisonSlideData) -> None:
        """Create the Plan Comparison slide"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        num_plans = len(data.plans)
        if num_plans == 0:
            return

        # Add title
        title_left = Inches(0.5)
        title_top = Inches(0.3)
        title_width = Inches(12.33)

        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, Inches(0.5))
        title_tf = title_box.text_frame
        title_tf.paragraphs[0].text = "Plan Benefit Comparison"
        title_tf.paragraphs[0].font.name = 'Poppins'
        title_tf.paragraphs[0].font.size = Pt(24)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = COLORS['black']

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(title_left, Inches(0.75), title_width, Inches(0.3))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.paragraphs[0].text = "Current employer plan vs. marketplace alternatives"
        subtitle_tf.paragraphs[0].font.name = 'Poppins'
        subtitle_tf.paragraphs[0].font.size = Pt(12)
        subtitle_tf.paragraphs[0].font.color.rgb = COLORS['detail_text']

        # Define rows for the comparison table
        rows = [
            ('header', 'Plan'),
            ('data', 'Age 21 Premium', 'age_21_premium'),
            ('data', 'Total Monthly Premium*', 'total_premium'),
            ('section', 'PLAN FEATURES'),
            ('data', 'Plan Type', 'plan_type'),
            ('data', 'HSA Eligible', 'hsa_eligible'),
            ('data', 'Coinsurance', 'coinsurance_pct'),
            ('data', 'Deductible (Individual)', 'individual_deductible'),
            ('data', 'Deductible (Family)', 'family_deductible'),
            ('data', 'OOP Max (Individual)', 'individual_oop_max'),
            ('data', 'OOP Max (Family)', 'family_oop_max'),
        ]

        # Create table
        table_left = Inches(0.4)
        table_top = Inches(1.15)
        table_width = Inches(12.5)
        num_rows = len(rows)
        row_height = Inches(0.38)  # Increased for content
        header_height = Inches(0.95)  # Increased for plan name + issuer + AV
        table_height = header_height + (num_rows - 1) * row_height

        table = slide.shapes.add_table(num_rows, num_plans + 1,
                                        table_left, table_top,
                                        table_width, table_height).table

        # Set column widths
        label_width = Inches(2.0)
        plan_width = (table_width - label_width) / num_plans
        table.columns[0].width = int(label_width)
        for i in range(1, num_plans + 1):
            table.columns[i].width = int(plan_width)

        # Set row heights
        table.rows[0].height = int(header_height)
        for i in range(1, num_rows):
            table.rows[i].height = int(row_height)

        # Build header row (row 0)
        cell = table.cell(0, 0)
        self._set_cell_fill(cell, COLORS['row_label_bg'])
        self._set_cell_text(cell, "", COLORS['row_label_text'])

        for col_idx, plan in enumerate(data.plans, start=1):
            cell = table.cell(0, col_idx)

            if plan.is_current:
                self._set_cell_fill(cell, COLORS['current_bg'])
                header_color = COLORS['current_header']
            else:
                self._set_cell_fill(cell, COLORS['white'])
                header_color = get_metal_color(plan.metal_level)

            # Plan name
            self._set_cell_text(cell, plan.plan_name, COLORS['black'],
                               font_size=10, bold=True)

            # Issuer name
            if plan.issuer_name:
                self._add_cell_line(cell, plan.issuer_name, COLORS['detail_text'], font_size=9)

            # AV% if available
            if plan.actuarial_value and not plan.is_current:
                self._add_cell_line(cell, f"{plan.actuarial_value:.0f}% AV", COLORS['detail_text'], font_size=9)

            # Metal chip text
            if plan.is_current:
                # plan_type is now a tuple (value, comparison)
                plan_type_val = plan.plan_type[0] if isinstance(plan.plan_type, tuple) else plan.plan_type
                chip_text = plan_type_val or "Group"
            else:
                chip_text = plan.metal_level.upper()
                # hsa_eligible is a tuple (value, comparison)
                hsa_val = plan.hsa_eligible[0] if isinstance(plan.hsa_eligible, tuple) else plan.hsa_eligible
                if hsa_val:
                    chip_text += " - HSA"

            self._add_cell_line(cell, chip_text, header_color, font_size=9, bold=True)

        # Build data rows
        for row_idx, row_def in enumerate(rows):
            if row_idx == 0:
                continue  # Skip header, already handled

            row_type = row_def[0]
            label = row_def[1]

            # Label cell
            cell = table.cell(row_idx, 0)

            if row_type == 'section':
                # Section header - spans visual distinction
                self._set_cell_fill(cell, COLORS['section_bg'])
                self._set_cell_text(cell, label, COLORS['section_text'],
                                   font_size=8, bold=False, align=PP_ALIGN.LEFT)
                # Fill rest of row with section bg
                for col_idx in range(1, num_plans + 1):
                    c = table.cell(row_idx, col_idx)
                    self._set_cell_fill(c, COLORS['section_bg'])
            else:
                # Data row
                attr_name = row_def[2] if len(row_def) > 2 else None

                self._set_cell_fill(cell, COLORS['white'])

                # Special handling for Total Premium label with employee info
                if attr_name == 'total_premium' and data.employee_count > 0:
                    self._set_cell_text(cell, label, COLORS['row_label_text'],
                                       font_size=10, bold=True, align=PP_ALIGN.LEFT)
                    age_text = f", avg age {data.avg_age:.0f}" if data.avg_age > 0 else ""
                    self._add_cell_line(cell, f"{data.employee_count} employees{age_text}",
                                       COLORS['detail_text'], font_size=8, align=PP_ALIGN.LEFT)
                else:
                    self._set_cell_text(cell, label, COLORS['row_label_text'],
                                       font_size=10, bold=False, align=PP_ALIGN.LEFT)

                # Data cells
                for col_idx, plan in enumerate(data.plans, start=1):
                    cell = table.cell(row_idx, col_idx)

                    if plan.is_current:
                        self._set_cell_fill(cell, COLORS['current_bg'])
                        text_color = COLORS['black']
                    else:
                        # Get comparison result using same logic as UI:
                        # - Premium rows: use dedicated comparison fields
                        # - Benefit rows: read from tuple
                        comparison = None

                        if attr_name == 'age_21_premium':
                            # Use pre-calculated comparison (same logic as UI)
                            comparison = plan.age_21_premium_comparison
                        elif attr_name == 'total_premium':
                            # Use pre-calculated comparison (same logic as UI)
                            comparison = plan.total_premium_comparison
                        elif attr_name and hasattr(plan, attr_name):
                            # Benefit rows: comparison stored in tuple
                            attr_val = getattr(plan, attr_name)
                            if isinstance(attr_val, tuple) and len(attr_val) == 2:
                                _, comparison = attr_val

                        # Apply color based on comparison result
                        if comparison == 'better':
                            self._set_cell_fill(cell, COLORS['better_bg'])
                        elif comparison == 'worse':
                            self._set_cell_fill(cell, COLORS['worse_bg'])
                        elif comparison == 'similar':
                            self._set_cell_fill(cell, COLORS['similar_bg'])
                        else:
                            self._set_cell_fill(cell, COLORS['white'])

                        text_color = COLORS['black']

                    # Format and display value
                    display_value = self._get_display_value(plan, attr_name)
                    self._set_cell_text(cell, display_value, text_color, font_size=10)

                    # Add detail line for specific rows
                    detail = self._get_detail_value(plan, attr_name, data)
                    if detail:
                        self._add_cell_line(cell, detail, COLORS['detail_text'], font_size=8)

        # Add footnote - position below table with adequate spacing
        if data.footnote:
            footnote_top = table_top + table_height + Inches(0.25)
            footnote_box = slide.shapes.add_textbox(table_left, footnote_top,
                                                     table_width, Inches(0.5))
            footnote_tf = footnote_box.text_frame
            footnote_tf.word_wrap = True
            footnote_tf.paragraphs[0].text = data.footnote
            footnote_tf.paragraphs[0].font.name = 'Poppins'
            footnote_tf.paragraphs[0].font.size = Pt(8)
            footnote_tf.paragraphs[0].font.color.rgb = COLORS['detail_text']

        # Add legend - position at bottom
        legend_top = Inches(7.0)
        legend_box = slide.shapes.add_textbox(table_left, legend_top, Inches(10), Inches(0.3))
        legend_tf = legend_box.text_frame
        legend_tf.paragraphs[0].text = "Green = Equivalent or better  |  Blue = Similar (within 5%)  |  Red = Less generous"
        legend_tf.paragraphs[0].font.name = 'Poppins'
        legend_tf.paragraphs[0].font.size = Pt(9)
        legend_tf.paragraphs[0].font.color.rgb = COLORS['detail_text']

        # Add decorative images
        self._add_decoratives(slide)

    def _get_display_value(self, plan: PlanColumnData, attr_name: str) -> str:
        """Get formatted display value for an attribute"""
        if attr_name is None:
            return ""

        if attr_name == 'age_21_premium':
            if plan.is_current:
                # Show current plan premium (EE-only, all ages) - not age-banded like marketplace
                display_prem = plan.renewal_age_21_premium or plan.current_age_21_premium
                if display_prem:
                    return self._format_currency(display_prem)
                return "—"
            return self._format_currency(plan.age_21_premium)

        elif attr_name == 'total_premium':
            if plan.is_current:
                # Show current/renewal for current plan
                curr = plan.current_total_premium
                renew = plan.renewal_total_premium
                if curr and renew and curr != renew:
                    return f"${curr:,.0f} / ${renew:,.0f}"
                elif renew:
                    return f"${renew:,.0f}"
                elif curr:
                    return f"${curr:,.0f}"
                return "—"
            return self._format_currency(plan.total_premium)

        elif attr_name == 'plan_type':
            if isinstance(plan.plan_type, tuple):
                return plan.plan_type[0] or "—"
            return plan.plan_type or "—"

        elif attr_name == 'hsa_eligible':
            if isinstance(plan.hsa_eligible, tuple):
                return "Yes" if plan.hsa_eligible[0] else "No"
            return "Yes" if plan.hsa_eligible else "No"

        elif attr_name == 'coinsurance_pct':
            if isinstance(plan.coinsurance_pct, tuple):
                return f"{plan.coinsurance_pct[0]}%"
            return f"{plan.coinsurance_pct}%"

        elif attr_name in ['individual_deductible', 'family_deductible',
                           'individual_oop_max', 'family_oop_max']:
            attr_val = getattr(plan, attr_name, (None, 'similar'))
            if isinstance(attr_val, tuple):
                value = attr_val[0]
            else:
                value = attr_val
            return self._format_currency(value)

        return "—"

    def _get_detail_value(self, plan: PlanColumnData, attr_name: str, data: PlanComparisonSlideData) -> str:
        """Get detail text for specific rows"""
        if attr_name == 'age_21_premium':
            if plan.is_current:
                # Note that group plans use flat rates, not age-banded
                display_prem = plan.renewal_age_21_premium or plan.current_age_21_premium
                if display_prem:
                    return "(EE-only, all ages)"
            return ""

        if attr_name == 'total_premium':
            if plan.is_current:
                # Show current/renewal averages for current plan
                curr_avg = plan.current_avg_premium
                renew_avg = plan.renewal_avg_premium
                if curr_avg and renew_avg and curr_avg != renew_avg:
                    return f"${curr_avg:,.0f} / ${renew_avg:,.0f} avg"
                elif renew_avg:
                    return f"${renew_avg:,.0f} avg"
                elif curr_avg:
                    return f"${curr_avg:,.0f} avg"
            else:
                # Marketplace plan - show diff from renewal and avg
                parts = []
                # Get renewal total from current plan in data
                renewal_total = None
                for p in data.plans:
                    if p.is_current and p.renewal_total_premium:
                        renewal_total = p.renewal_total_premium
                        break
                if renewal_total and plan.total_premium:
                    diff = plan.total_premium - renewal_total
                    if diff < 0:
                        parts.append(f"${abs(diff):,.0f} less")
                    elif diff > 0:
                        parts.append(f"${diff:,.0f} more")
                if plan.avg_premium:
                    parts.append(f"${plan.avg_premium:,.0f} avg")
                return "\n".join(parts) if parts else ""

        if attr_name == 'age_21_premium':
            if plan.is_current:
                # No detail for current plan - Age 21 Premium is N/A for group plans
                return ""
            else:
                # Marketplace plan - show diff from renewal premium
                renewal_age_21 = None
                for p in data.plans:
                    if p.is_current and p.renewal_age_21_premium:
                        renewal_age_21 = p.renewal_age_21_premium
                        break
                if renewal_age_21 and plan.age_21_premium:
                    diff = plan.age_21_premium - renewal_age_21
                    if diff < 0:
                        return f"${abs(diff):,.0f} less"
                    elif diff > 0:
                        return f"${diff:,.0f} more"

        return ""

    def _add_decoratives(self, slide) -> None:
        """Add decorative corner and edge images to the slide"""
        corner_size = Inches(1.2)
        if CORNER_IMAGE.exists():
            slide.shapes.add_picture(
                str(CORNER_IMAGE),
                left=Inches(0),
                top=Inches(0),
                width=corner_size,
                height=corner_size
            )

        edge_width = Inches(2.8)
        edge_height = Inches(0.7)
        if EDGE_IMAGE.exists():
            slide.shapes.add_picture(
                str(EDGE_IMAGE),
                left=SLIDE_WIDTH - edge_width,
                top=SLIDE_HEIGHT - edge_height,
                width=edge_width,
                height=edge_height
            )

    def generate(self) -> BytesIO:
        """Generate PowerPoint and return as BytesIO buffer"""
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer

    def save(self, output_path: str) -> None:
        """Save PowerPoint to file"""
        self.prs.save(output_path)


def generate_plan_comparison_slide(data: PlanComparisonSlideData) -> BytesIO:
    """
    Convenience function to generate a Plan Comparison slide.

    Args:
        data: PlanComparisonSlideData with values to populate

    Returns:
        BytesIO buffer with generated PowerPoint
    """
    generator = PlanComparisonSlideGenerator()
    generator.create_slide(data)
    return generator.generate()


if __name__ == "__main__":
    # Test with sample data
    print("Testing Plan Comparison Slide Generator...")

    data = PlanComparisonSlideData(
        employee_count=42,
        avg_age=38.5,
        footnote="*Assumes all 42 employees are in Rating Area 12.",
        plans=[
            PlanColumnData(
                plan_name="Blue Preferred Plus POS HSA",
                issuer_name="BCBS",
                plan_type=("POS", 'similar'),
                is_current=True,
                hsa_eligible=(True, 'similar'),
                total_premium=36554,
                avg_premium=870,
                individual_deductible=(2650, 'similar'),
                individual_oop_max=(6650, 'similar'),
                coinsurance_pct=(0, 'similar'),
            ),
            PlanColumnData(
                plan_name="Anthem Silver Pathway",
                issuer_name="Anthem Blue Cross",
                plan_type=("HMO", 'similar'),
                metal_level="Silver",
                actuarial_value=70,
                hsa_eligible=(False, 'similar'),
                age_21_premium=471,
                age_21_premium_comparison='worse',
                total_premium=42020,
                total_premium_comparison='worse',
                avg_premium=1001,
                affordable_contribution=17612,
                contribution_range="$0-$955/ee",
                individual_deductible=(6000, 'worse'),
                individual_oop_max=(9450, 'worse'),
                coinsurance_pct=(20, 'worse'),
            ),
            PlanColumnData(
                plan_name="Dean Gold HSA",
                issuer_name="Dean Health Plan",
                plan_type=("HMO", 'similar'),
                metal_level="Gold",
                actuarial_value=78,
                hsa_eligible=(True, 'better'),
                age_21_premium=514,
                age_21_premium_comparison='worse',
                total_premium=45831,
                total_premium_comparison='worse',
                avg_premium=1091,
                affordable_contribution=20912,
                contribution_range="$0-$1,083/ee",
                individual_deductible=(2400, 'better'),
                individual_oop_max=(5500, 'better'),
                coinsurance_pct=(20, 'worse'),
            ),
        ],
    )

    generator = PlanComparisonSlideGenerator()
    generator.create_slide(data)

    output_path = Path(__file__).parent / 'test_plan_comparison.pptx'
    generator.save(str(output_path))

    print(f"Generated test slide: {output_path}")
    print("Test complete!")
