"""
Cooperative Health Plan Comparison Slide Generator

Generates a single PowerPoint slide with the Cooperative Health Plan Comparison table
that matches the design from the ICHRA dashboard.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from typing import Dict, List, Optional, Any
from io import BytesIO
from pathlib import Path
from dataclasses import dataclass, field
import pandas as pd


# Color scheme matching Figma design
COLORS = {
    # Current 2025 - Blue
    'current_bg': RGBColor(0xEF, 0xF6, 0xFF),      # #EFF6FF
    'current_text': RGBColor(0x1C, 0x39, 0x8E),    # #1C398E
    'current_header': RGBColor(0xDB, 0xEA, 0xFE),  # #DBEAFE

    # Renewal 2026 - Red
    'renewal_bg': RGBColor(0xFE, 0xF2, 0xF2),      # #FEF2F2
    'renewal_text': RGBColor(0x82, 0x18, 0x1A),    # #82181A
    'renewal_header': RGBColor(0xFE, 0xCA, 0xCA),  # #FECACA

    # HAS - Green
    'has_bg': RGBColor(0xF0, 0xFD, 0xF4),          # #F0FDF4
    'has_text': RGBColor(0x0D, 0x54, 0x2B),        # #0D542B
    'has_header': RGBColor(0xBB, 0xF7, 0xD0),      # #BBF7D0

    # Sedera - Blue
    'sedera_bg': RGBColor(0xEF, 0xF6, 0xFF),       # #EFF6FF
    'sedera_text': RGBColor(0x1C, 0x39, 0x8E),     # #1C398E
    'sedera_header': RGBColor(0xDB, 0xEA, 0xFE),   # #DBEAFE

    # Neutral
    'row_label_bg': RGBColor(0xF9, 0xFA, 0xFB),    # #F9FAFB
    'row_label_text': RGBColor(0x37, 0x41, 0x51),  # #374151
    'total_bg': RGBColor(0xF3, 0xF4, 0xF6),        # #F3F4F6
    'total_text': RGBColor(0x11, 0x18, 0x27),      # #111827
    'gap_text': RGBColor(0x9C, 0xA3, 0xAF),        # #9CA3AF - lighter gray for breakdown
    'border': RGBColor(0xE5, 0xE7, 0xEB),          # #E5E7EB
    'white': RGBColor(0xFF, 0xFF, 0xFF),           # #FFFFFF
}

# Slide dimensions (standard widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Decorative image paths
DECORATIVES_DIR = Path(__file__).parent / "decoratives"
EDGE_IMAGE = DECORATIVES_DIR / "glove-tile-edge-h-fade.png"

# Banner image path
BANNER_IMAGE = DECORATIVES_DIR / "PPT_header.png"


@dataclass
class PlanColumnData:
    """Data for a single plan column (HAS or Sedera with specific IUA/deductible)"""
    key: str  # Unique key (e.g., "HAS $1k", "Sedera $500")
    label: str  # Display label (e.g., "HAS $1k")
    subtitle: str  # Subtitle (e.g., "Health Access Solutions")
    plan_type: str  # 'has' or 'sedera'
    total: float = 0.0  # Monthly total for this tier
    min_rate: float = 0.0  # Min per-employee rate
    max_rate: float = 0.0  # Max per-employee rate


@dataclass
class TierData:
    """Data for a single coverage tier row"""
    name: str  # Display name (e.g., "Employee Only")
    code: str  # Tier code (e.g., "EE")
    current_total: float = 0.0
    current_base: float = 0.0
    current_gap: float = 0.0
    renewal_total: float = 0.0
    renewal_base: float = 0.0
    renewal_gap: float = 0.0
    # Dynamic plan columns - dict keyed by plan column key
    plan_columns: Dict[str, PlanColumnData] = field(default_factory=dict)
    # Employee count and per-employee rates
    employee_count: int = 0
    avg_age: int = 0  # Average age of employees in this tier
    current_rate_per_ee: float = 0.0  # Average per-employee base rate (current)
    renewal_rate_per_ee: float = 0.0  # Average per-employee base rate (renewal)
    gap_rate_per_ee: float = 0.0      # Fixed gap rate per employee for this tier


@dataclass
class PlanColumnTotals:
    """Totals for a single plan column across all tiers"""
    key: str
    label: str
    subtitle: str
    plan_type: str  # 'has' or 'sedera'
    monthly_total: float = 0.0
    annual_total: float = 0.0
    savings_amount: float = 0.0
    savings_pct: float = 0.0
    # Admin fee for this plan column
    admin_fee_pepm: float = 0.0  # Per-employee-per-month rate
    admin_fee_total: float = 0.0  # Total monthly admin fee


@dataclass
class CooperativeHealthData:
    """Data container for Cooperative Health Plan Comparison slide"""

    # Tier breakdown
    tiers: List[TierData] = field(default_factory=list)

    # Dynamic plan columns (ordered list)
    plan_columns: List[PlanColumnTotals] = field(default_factory=list)

    # Monthly totals for fixed columns
    total_current: float = 0.0
    total_renewal: float = 0.0

    # Annual totals for fixed columns
    annual_current: float = 0.0
    annual_renewal: float = 0.0

    # Gap insurance details
    has_gap: bool = False
    total_gap_monthly: float = 0.0
    gap_rate_ee: float = 0.0
    gap_rate_es: float = 0.0
    gap_rate_ec: float = 0.0
    gap_rate_f: float = 0.0

    # Preventive care details
    has_preventive: bool = False
    preventive_total: float = 0.0  # Total monthly preventive care cost
    preventive_rate: float = 0.0  # Per-employee rate
    preventive_employee_count: int = 0
    preventive_include_dpc: bool = True  # "with DPC" vs "without DPC"
    # Preventive by tier: {tier_code: {'rate': X, 'count': Y, 'total': Z}}
    preventive_by_tier: Dict[str, Dict] = field(default_factory=dict)

    # Admin fee flag (True if any plan column has admin fee > 0)
    has_admin_fees: bool = False

    # Renewal data availability - when False, hide renewal column and compare to current
    has_renewal_data: bool = True
    comparison_label: str = "Savings vs Renewal"  # "Savings vs Renewal" or "Savings vs Current"

    # Client info
    client_name: str = ""

    @classmethod
    def from_session_state(cls, session_state) -> 'CooperativeHealthData':
        """
        Build CooperativeHealthData from Streamlit session_state.

        Args:
            session_state: st.session_state object

        Returns:
            Populated CooperativeHealthData instance
        """
        data = cls()

        census_df = session_state.get('census_df')
        financial_summary = session_state.get('financial_summary', {})
        client_name = session_state.get('client_name', '')

        data.client_name = client_name

        if census_df is None or census_df.empty:
            return data

        # Get multi-metal scenario data if available
        multi_metal = financial_summary.get('multi_metal_scenario', {})
        tier_totals = multi_metal.get('tier_totals', {})
        hap_analysis = multi_metal.get('hap_analysis', {})

        # Check for gap insurance data
        has_gap_col = 'gap_insurance_monthly' in census_df.columns
        if has_gap_col:
            gap_values = census_df['gap_insurance_monthly'].fillna(0)
            if gap_values.sum() > 0:
                data.has_gap = True
                data.total_gap_monthly = float(gap_values.sum())

                # Get gap rates by tier (fixed per tier - take first non-zero)
                for code, attr in [('EE', 'gap_rate_ee'), ('ES', 'gap_rate_es'),
                                   ('EC', 'gap_rate_ec'), ('F', 'gap_rate_f')]:
                    tier_emps = census_df[census_df['family_status'] == code]
                    if not tier_emps.empty:
                        gap_vals = tier_emps['gap_insurance_monthly'].dropna()
                        gap_vals = gap_vals[gap_vals > 0]
                        if len(gap_vals) > 0:
                            setattr(data, attr, float(gap_vals.iloc[0]))

        # Define tier info
        tier_info = {
            'Employee Only': {'code': 'EE', 'display': 'Employee Only'},
            'Employee + Spouse': {'code': 'ES', 'display': 'Employee + Spouse'},
            'Employee + Children': {'code': 'EC', 'display': 'Employee + Children'},
            'Family': {'code': 'F', 'display': 'Family'},
        }

        # Build tier data
        for tier_name, info in tier_info.items():
            tier_data = tier_totals.get(tier_name, {})
            code = info['code']

            # Calculate gap for this tier
            tier_gap = 0.0
            if data.has_gap:
                tier_emps = census_df[census_df['family_status'] == code]
                if not tier_emps.empty:
                    tier_gap = float(tier_emps['gap_insurance_monthly'].fillna(0).sum())

            # Get HAS analysis for this tier
            hap_1k = hap_analysis.get('hap_1k', {}).get(tier_name, {})
            hap_25k = hap_analysis.get('hap_2.5k', {}).get(tier_name, {})

            # Base premiums (without gap)
            current_base = tier_data.get('current_total', 0) or 0
            renewal_base = tier_data.get('renewal_2026', 0) or 0

            # Total includes gap insurance if present
            current_total = current_base + tier_gap
            renewal_total = renewal_base + tier_gap

            # Get employee count for this tier
            tier_emps = census_df[census_df['family_status'] == code]
            employee_count = len(tier_emps)

            # Get typical per-employee current rate (most common value from census)
            if 'current_ee_monthly' in census_df.columns and 'current_er_monthly' in census_df.columns:
                ee_vals = tier_emps['current_ee_monthly'].fillna(0)
                er_vals = tier_emps['current_er_monthly'].fillna(0)
                per_ee_rates = ee_vals + er_vals
                non_zero_rates = per_ee_rates[per_ee_rates > 0]
                if len(non_zero_rates) > 0:
                    try:
                        current_rate_per_ee = float(non_zero_rates.mode().iloc[0])
                    except (IndexError, ValueError):
                        current_rate_per_ee = float(non_zero_rates.iloc[0])
                else:
                    current_rate_per_ee = 0
            else:
                current_rate_per_ee = 0

            # Get typical per-employee renewal rate (most common value from census)
            if 'projected_2026_premium' in census_df.columns:
                renewal_vals = tier_emps['projected_2026_premium'].fillna(0)
                non_zero_renewal = renewal_vals[renewal_vals > 0]
                if len(non_zero_renewal) > 0:
                    try:
                        renewal_rate_per_ee = float(non_zero_renewal.mode().iloc[0])
                    except (IndexError, ValueError):
                        renewal_rate_per_ee = float(non_zero_renewal.iloc[0])
                else:
                    renewal_rate_per_ee = 0
            else:
                renewal_rate_per_ee = 0

            # Get gap rate for this tier (fixed per tier)
            gap_rate_per_ee = getattr(data, f'gap_rate_{code.lower()}', 0)

            # Calculate average age for this tier
            avg_age = 0
            if 'age' in tier_emps.columns and employee_count > 0:
                avg_age = round(tier_emps['age'].mean())

            tier = TierData(
                name=tier_name,
                code=code,
                current_total=current_total,
                current_base=current_base,
                current_gap=tier_gap,
                renewal_total=renewal_total,
                renewal_base=renewal_base,
                renewal_gap=tier_gap,
                hap_1k_total=hap_1k.get('total', 0) or 0,
                hap_1k_min=hap_1k.get('min_rate', 0) or 0,
                hap_1k_max=hap_1k.get('max_rate', 0) or 0,
                hap_25k_total=hap_25k.get('total', 0) or 0,
                hap_25k_min=hap_25k.get('min_rate', 0) or 0,
                hap_25k_max=hap_25k.get('max_rate', 0) or 0,
                employee_count=employee_count,
                avg_age=avg_age,
                current_rate_per_ee=current_rate_per_ee,
                renewal_rate_per_ee=renewal_rate_per_ee,
                gap_rate_per_ee=gap_rate_per_ee,
            )
            data.tiers.append(tier)

        # Calculate totals
        data.total_current = sum(t.current_total for t in data.tiers)
        data.total_renewal = sum(t.renewal_total for t in data.tiers)
        data.total_hap_1k = sum(t.hap_1k_total for t in data.tiers)
        data.total_hap_25k = sum(t.hap_25k_total for t in data.tiers)

        return data


class CooperativeHealthSlideGenerator:
    """Generate Cooperative Health Plan Comparison PowerPoint slide"""

    def __init__(self):
        """Initialize with a new blank presentation"""
        self.prs = Presentation()
        # Set slide dimensions to standard widescreen
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    def _set_cell_fill(self, cell, color: RGBColor):
        """Set cell background color"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    def _set_cell_text(self, cell, text: str, color: RGBColor,
                       font_size: int = 11, bold: bool = False,
                       align: PP_ALIGN = PP_ALIGN.CENTER):
        """Set cell text with formatting"""
        cell.text = text
        cell.text_frame.paragraphs[0].alignment = align
        cell.text_frame.paragraphs[0].font.name = 'Poppins'
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.text_frame.paragraphs[0].font.color.rgb = color
        cell.text_frame.paragraphs[0].font.bold = bold
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _add_cell_line(self, cell, text: str, color: RGBColor,
                       font_size: int = 10, bold: bool = False,
                       align: PP_ALIGN = PP_ALIGN.CENTER):
        """Add additional line to cell"""
        p = cell.text_frame.add_paragraph()
        p.text = text
        p.alignment = align
        p.font.name = 'Poppins'
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold

    def _format_currency(self, value: float) -> str:
        """Format value as currency"""
        return f"${value:,.0f}"

    def create_slide(self, data: CooperativeHealthData) -> None:
        """
        Create the Cooperative Health Plan Comparison slide.

        Args:
            data: CooperativeHealthData with values to display
        """
        # Add blank slide
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)

        # Add banner at top
        if BANNER_IMAGE.exists():
            slide.shapes.add_picture(
                str(BANNER_IMAGE),
                left=Inches(0),
                top=Inches(0),
                width=SLIDE_WIDTH,
                height=Inches(0.25)
            )

        # Banner offset for content positioning
        banner_offset = Inches(0.2)

        # Add company name (top right)
        if data.client_name:
            company_box = slide.shapes.add_textbox(
                Inches(9.5), Inches(0.3), Inches(3.5), Inches(0.3)
            )
            company_tf = company_box.text_frame
            company_tf.paragraphs[0].text = data.client_name
            company_tf.paragraphs[0].font.name = 'Poppins'
            company_tf.paragraphs[0].font.size = Pt(14)
            company_tf.paragraphs[0].font.bold = True
            company_tf.paragraphs[0].font.color.rgb = COLORS['total_text']
            company_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Add title
        title_left = Inches(0.5)
        title_top = Inches(0.35) + banner_offset
        title_width = Inches(9)
        title_height = Inches(0.6)

        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_tf = title_box.text_frame
        title_tf.paragraphs[0].text = "Cooperative Health Plan Comparison"
        title_tf.paragraphs[0].font.name = 'Poppins'
        title_tf.paragraphs[0].font.size = Pt(26)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = COLORS['total_text']

        # Add subtitle/description
        subtitle_top = Inches(0.85) + banner_offset
        subtitle_height = Inches(0.4)

        subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top, Inches(12.33), subtitle_height)
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.paragraphs[0].text = "Traditional group plan vs. cooperative health plan alternatives"
        subtitle_tf.paragraphs[0].font.name = 'Poppins'
        subtitle_tf.paragraphs[0].font.size = Pt(14)
        subtitle_tf.paragraphs[0].font.color.rgb = COLORS['gap_text']

        # Dynamic column count: fixed columns + plan columns
        # 3 fixed (Coverage Type, Current, Renewal) when has_renewal_data
        # 2 fixed (Coverage Type, Current) when no renewal data
        num_plan_cols = len(data.plan_columns)
        num_fixed_cols = 3 if data.has_renewal_data else 2
        num_cols = num_fixed_cols + num_plan_cols

        # Create table
        table_left = Inches(0.5)
        table_top = Inches(1.4) + banner_offset
        table_width = Inches(12.33)

        # Calculate row heights
        header_height = Inches(0.5)
        # Increase tier height if gap or preventive data exists (extra lines)
        tier_height = Inches(0.75) if (data.has_gap or data.has_preventive) else Inches(0.5)
        summary_height = Inches(0.45)

        # Calculate dynamic row count
        # Base: Header + 4 tiers + Total Monthly + Annual Total + Savings vs Renewal = 8
        num_rows = 8
        if data.has_admin_fees:
            num_rows += 1  # Admin Fee row (Grand Total merged into Total Monthly)

        num_summary_rows = 3 + (1 if data.has_admin_fees else 0)
        table_height = header_height + (4 * tier_height) + (num_summary_rows * summary_height)

        # Add table with dynamic columns
        table = slide.shapes.add_table(num_rows, num_cols, table_left, table_top,
                                        table_width, table_height).table

        # Calculate column widths dynamically
        first_col_width = Inches(2.5)
        remaining_width = 12.33 - 2.5  # Total width minus first column
        other_col_width = Inches(remaining_width / (num_cols - 1))

        table.columns[0].width = first_col_width
        for i in range(1, num_cols):
            table.columns[i].width = other_col_width

        # Build dynamic headers and colors - conditionally include renewal column
        if data.has_renewal_data:
            headers = ['Coverage Type', 'Current 2025', 'Renewal 2026']
            header_colors = [
                (COLORS['row_label_bg'], COLORS['row_label_text']),
                (COLORS['current_header'], COLORS['current_text']),
                (COLORS['renewal_header'], COLORS['renewal_text']),
            ]
            data_text_colors = [
                COLORS['row_label_text'],
                COLORS['current_text'],
                COLORS['renewal_text'],
            ]
        else:
            headers = ['Coverage Type', 'Current 2025']
            header_colors = [
                (COLORS['row_label_bg'], COLORS['row_label_text']),
                (COLORS['current_header'], COLORS['current_text']),
            ]
            data_text_colors = [
                COLORS['row_label_text'],
                COLORS['current_text'],
            ]

        # Add plan column headers and colors
        for plan_col in data.plan_columns:
            headers.append(plan_col.label)
            if plan_col.plan_type == 'sedera':
                header_colors.append((COLORS['sedera_header'], COLORS['sedera_text']))
                data_text_colors.append(COLORS['sedera_text'])
            else:  # 'has'
                header_colors.append((COLORS['has_header'], COLORS['has_text']))
                data_text_colors.append(COLORS['has_text'])

        # Style header row
        for col_idx, (header_text, (bg_color, text_color)) in enumerate(zip(headers, header_colors)):
            cell = table.cell(0, col_idx)
            self._set_cell_fill(cell, bg_color)
            self._set_cell_text(cell, header_text, text_color, font_size=12, bold=True)

        # Add "with Preventive" subtitle to HAS column headers when preventive is enabled
        if data.has_preventive:
            for col_offset, plan_col in enumerate(data.plan_columns):
                if plan_col.plan_type == 'has':
                    col_idx = num_fixed_cols + col_offset
                    cell = table.cell(0, col_idx)
                    self._add_cell_line(cell, "with Preventive", COLORS['has_text'], font_size=9)

        # Style tier rows (white background for data cells)
        for row_idx, tier in enumerate(data.tiers, start=1):
            # Coverage Type column - with employee count (left-aligned)
            cell = table.cell(row_idx, 0)
            self._set_cell_fill(cell, COLORS['row_label_bg'])
            self._set_cell_text(cell, tier.name, data_text_colors[0],
                               font_size=11, bold=True, align=PP_ALIGN.LEFT)
            if tier.employee_count > 0:
                age_suffix = f" · avg age {tier.avg_age}" if tier.avg_age > 0 else ""
                self._add_cell_line(cell, f"{tier.employee_count} employees{age_suffix}", COLORS['gap_text'],
                                   font_size=9, align=PP_ALIGN.LEFT)

            # Current 2025 column
            cell = table.cell(row_idx, 1)
            self._set_cell_fill(cell, COLORS['white'])
            self._set_cell_text(cell, self._format_currency(tier.current_total),
                               data_text_colors[1], font_size=12, bold=True)
            if tier.current_rate_per_ee > 0:
                if data.has_gap and tier.gap_rate_per_ee > 0:
                    breakdown = f"${tier.current_rate_per_ee:,.0f} + ${tier.gap_rate_per_ee:,.0f} gap"
                else:
                    breakdown = f"${tier.current_rate_per_ee:,.0f}"
                self._add_cell_line(cell, breakdown, COLORS['gap_text'], font_size=9)

            # Renewal 2026 column (only if has_renewal_data)
            if data.has_renewal_data:
                cell = table.cell(row_idx, 2)
                self._set_cell_fill(cell, COLORS['white'])
                self._set_cell_text(cell, self._format_currency(tier.renewal_total),
                                   data_text_colors[2], font_size=12, bold=True)
                if tier.renewal_rate_per_ee > 0:
                    if data.has_gap and tier.gap_rate_per_ee > 0:
                        breakdown = f"${tier.renewal_rate_per_ee:,.0f} + ${tier.gap_rate_per_ee:,.0f} gap"
                    else:
                        breakdown = f"${tier.renewal_rate_per_ee:,.0f}"
                    self._add_cell_line(cell, breakdown, COLORS['gap_text'], font_size=9)

            # Dynamic plan columns
            for col_offset, plan_col in enumerate(data.plan_columns):
                col_idx = num_fixed_cols + col_offset
                cell = table.cell(row_idx, col_idx)
                self._set_cell_fill(cell, COLORS['white'])

                # Get tier data for this plan column
                tier_plan_data = tier.plan_columns.get(plan_col.key)
                if tier_plan_data:
                    self._set_cell_text(cell, self._format_currency(tier_plan_data.total),
                                       data_text_colors[col_idx], font_size=12, bold=True)
                    if tier_plan_data.min_rate > 0 and tier_plan_data.max_rate > 0:
                        rate_range = f"${tier_plan_data.min_rate:,.0f}-${tier_plan_data.max_rate:,.0f}"
                        self._add_cell_line(cell, rate_range, COLORS['gap_text'], font_size=9)
                    # Add preventive care breakdown if enabled for this tier (HAS only)
                    if data.has_preventive and plan_col.plan_type == 'has' and tier.code in data.preventive_by_tier:
                        prev_data = data.preventive_by_tier[tier.code]
                        if prev_data.get('count', 0) > 0:
                            prev_line = f"${prev_data['rate']:,.0f}/mo × {prev_data['count']} = +${prev_data['total']:,.0f}"
                            self._add_cell_line(cell, prev_line, COLORS['gap_text'], font_size=9)
                else:
                    self._set_cell_text(cell, "N/A", COLORS['gap_text'], font_size=12)

        # Track current row index (starts after tier rows)
        current_row = 5

        # Row: Admin Fee (if any plan has admin fee) - appears above Total Monthly
        if data.has_admin_fees:
            admin_row = current_row

            cell = table.cell(admin_row, 0)
            self._set_cell_fill(cell, COLORS['row_label_bg'])
            self._set_cell_text(cell, "Admin Fee", COLORS['row_label_text'],
                               font_size=11, bold=True, align=PP_ALIGN.LEFT)

            # Empty cell for Current
            cell = table.cell(admin_row, 1)
            self._set_cell_fill(cell, COLORS['white'])

            # Empty cell for Renewal (only if has_renewal_data)
            if data.has_renewal_data:
                cell = table.cell(admin_row, 2)
                self._set_cell_fill(cell, COLORS['white'])

            # Plan column admin fees
            for col_offset, plan_col in enumerate(data.plan_columns):
                col_idx = num_fixed_cols + col_offset
                cell = table.cell(admin_row, col_idx)
                self._set_cell_fill(cell, COLORS['white'])
                if plan_col.admin_fee_total > 0:
                    self._set_cell_text(cell, self._format_currency(plan_col.admin_fee_total),
                                       data_text_colors[col_idx], font_size=12, bold=True)
                    self._add_cell_line(cell, f"${plan_col.admin_fee_pepm:.2f} PEPM",
                                       COLORS['gap_text'], font_size=9)
                else:
                    self._set_cell_text(cell, "—", COLORS['gap_text'], font_size=12)

            current_row += 1

        # Row: Total Monthly (now includes admin fees - this is the grand total)
        total_monthly_row = current_row
        total_label = "Total Monthly*" if data.has_gap else "Total Monthly"

        cell = table.cell(total_monthly_row, 0)
        self._set_cell_fill(cell, COLORS['white'])
        self._set_cell_text(cell, total_label, COLORS['total_text'],
                           font_size=12, bold=True, align=PP_ALIGN.LEFT)

        # Current monthly total
        cell = table.cell(total_monthly_row, 1)
        self._set_cell_fill(cell, COLORS['white'])
        self._set_cell_text(cell, self._format_currency(data.total_current), data_text_colors[1],
                           font_size=13, bold=True)

        # Renewal monthly total (only if has_renewal_data)
        if data.has_renewal_data:
            cell = table.cell(total_monthly_row, 2)
            self._set_cell_fill(cell, COLORS['white'])
            self._set_cell_text(cell, self._format_currency(data.total_renewal), data_text_colors[2],
                               font_size=13, bold=True)

        # Plan column monthly totals (grand total = monthly_total + admin_fee_total)
        for col_offset, plan_col in enumerate(data.plan_columns):
            col_idx = num_fixed_cols + col_offset
            cell = table.cell(total_monthly_row, col_idx)
            self._set_cell_fill(cell, COLORS['white'])
            grand_total = plan_col.monthly_total + plan_col.admin_fee_total
            self._set_cell_text(cell, self._format_currency(grand_total),
                               data_text_colors[col_idx], font_size=13, bold=True)

        current_row += 1

        # Row: Annual Total
        annual_row = current_row
        cell = table.cell(annual_row, 0)
        self._set_cell_fill(cell, COLORS['white'])
        self._set_cell_text(cell, "Annual Total", COLORS['total_text'],
                           font_size=12, bold=True, align=PP_ALIGN.LEFT)

        # Current annual total
        cell = table.cell(annual_row, 1)
        self._set_cell_fill(cell, COLORS['white'])
        self._set_cell_text(cell, self._format_currency(data.annual_current), data_text_colors[1],
                           font_size=12, bold=True)

        # Renewal annual total (only if has_renewal_data)
        if data.has_renewal_data:
            cell = table.cell(annual_row, 2)
            self._set_cell_fill(cell, COLORS['white'])
            self._set_cell_text(cell, self._format_currency(data.annual_renewal), data_text_colors[2],
                               font_size=12, bold=True)

        # Plan column annual totals
        for col_offset, plan_col in enumerate(data.plan_columns):
            col_idx = num_fixed_cols + col_offset
            cell = table.cell(annual_row, col_idx)
            self._set_cell_fill(cell, COLORS['white'])
            self._set_cell_text(cell, self._format_currency(plan_col.annual_total),
                               data_text_colors[col_idx], font_size=12, bold=True)

        current_row += 1

        # Row: Savings (dynamic label based on comparison baseline)
        savings_row = current_row
        cell = table.cell(savings_row, 0)
        self._set_cell_fill(cell, COLORS['white'])
        self._set_cell_text(cell, data.comparison_label, COLORS['total_text'],
                           font_size=12, bold=True, align=PP_ALIGN.LEFT)

        # Baseline column (Current or Renewal): dash to indicate comparison base
        cell = table.cell(savings_row, 1)
        self._set_cell_fill(cell, COLORS['white'])
        if data.has_renewal_data:
            # Current column empty, Renewal column shows dash
            cell = table.cell(savings_row, 2)
            self._set_cell_fill(cell, COLORS['white'])
            self._set_cell_text(cell, "—", COLORS['gap_text'], font_size=12)
        else:
            # No renewal column, Current column shows dash
            self._set_cell_text(cell, "—", COLORS['gap_text'], font_size=12)

        # Plan column savings - dollar amount first, percentage on second line (matches UI)
        for col_offset, plan_col in enumerate(data.plan_columns):
            col_idx = num_fixed_cols + col_offset
            cell = table.cell(savings_row, col_idx)
            self._set_cell_fill(cell, COLORS['white'])

            # Determine color based on savings (green for positive, red for negative)
            if plan_col.savings_amount > 0:
                savings_color = RGBColor(0x00, 0xA6, 0x3E)  # Green #00a63e
            elif plan_col.savings_amount < 0:
                savings_color = RGBColor(0xDC, 0x26, 0x26)  # Red #dc2626
            else:
                savings_color = COLORS['gap_text']  # Gray for zero

            if plan_col.savings_amount != 0:
                # Dollar amount first - bold (matches UI)
                # Use absolute value for display - color indicates savings/cost increase
                self._set_cell_text(cell, self._format_currency(abs(plan_col.savings_amount)),
                                   savings_color, font_size=12, bold=True)
                # Percentage on second line - smaller (matches UI format)
                self._add_cell_line(cell, f"({plan_col.savings_pct:.0f}%)",
                                   savings_color, font_size=10, bold=False)
            else:
                self._set_cell_text(cell, "—", COLORS['gap_text'], font_size=12)

        # Add footnote if gap insurance exists
        if data.has_gap:
            footnote_top = table_top + table_height + Inches(0.15)
            footnote_box = slide.shapes.add_textbox(table_left, footnote_top,
                                                     table_width, Inches(0.3))
            footnote_tf = footnote_box.text_frame

            gap_breakdown = (f"* Includes {self._format_currency(data.total_gap_monthly)}/mo gap coverage "
                           f"(EE: {self._format_currency(data.gap_rate_ee)} | "
                           f"ES: {self._format_currency(data.gap_rate_es)} | "
                           f"EC: {self._format_currency(data.gap_rate_ec)} | "
                           f"F: {self._format_currency(data.gap_rate_f)})")

            footnote_tf.paragraphs[0].text = gap_breakdown
            footnote_tf.paragraphs[0].font.name = 'Poppins'
            footnote_tf.paragraphs[0].font.size = Pt(9)
            footnote_tf.paragraphs[0].font.color.rgb = COLORS['gap_text']

        # Build dynamic legend text
        legend_parts = []
        has_has = any(pc.plan_type == 'has' for pc in data.plan_columns)
        has_sedera = any(pc.plan_type == 'sedera' for pc in data.plan_columns)
        if has_has:
            legend_parts.append("HAS = Health Access Solutions with specified deductible")
        if has_sedera:
            legend_parts.append("Sedera = Health sharing with specified IUA")
        legend_parts.append("Rates shown are total monthly premiums")

        legend_top = Inches(6.8)
        legend_box = slide.shapes.add_textbox(table_left, legend_top, Inches(10), Inches(0.4))
        legend_tf = legend_box.text_frame
        legend_tf.paragraphs[0].text = " | ".join(legend_parts)
        legend_tf.paragraphs[0].font.name = 'Poppins'
        legend_tf.paragraphs[0].font.size = Pt(9)
        legend_tf.paragraphs[0].font.color.rgb = COLORS['gap_text']

        # Add decorative corner and edge images
        self._add_decoratives(slide)

    def _add_savings_hero(self, slide, data: CooperativeHealthData) -> None:
        """Add hero savings callout to bottom-right area"""
        # Find the plan column with the best savings
        if not data.plan_columns:
            return  # No plan columns to display

        best_plan = max(data.plan_columns, key=lambda p: p.savings_pct)
        if best_plan.savings_pct <= 0:
            return  # No savings to display

        # Determine colors based on plan type
        if best_plan.plan_type == 'sedera':
            bg_color = COLORS['sedera_header']
            text_color = COLORS['sedera_text']
        else:  # 'has'
            bg_color = COLORS['has_header']
            text_color = COLORS['has_text']

        # Position: bottom-right area, above the decorative edge
        hero_left = Inches(9.5)
        hero_top = Inches(5.2)
        hero_width = Inches(3.3)
        hero_height = Inches(1.4)

        # Add background shape (rounded rectangle effect via textbox with fill)
        from pptx.enum.shapes import MSO_SHAPE
        hero_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            hero_left, hero_top, hero_width, hero_height
        )
        hero_shape.fill.solid()
        hero_shape.fill.fore_color.rgb = bg_color
        hero_shape.line.fill.background()  # No border

        # Add "SAVE UP TO" label
        label_box = slide.shapes.add_textbox(hero_left, hero_top + Inches(0.15), hero_width, Inches(0.3))
        label_tf = label_box.text_frame
        label_tf.paragraphs[0].text = "SAVE UP TO"
        label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        label_tf.paragraphs[0].font.name = 'Poppins'
        label_tf.paragraphs[0].font.size = Pt(12)
        label_tf.paragraphs[0].font.bold = True
        label_tf.paragraphs[0].font.color.rgb = text_color

        # Add big percentage number
        pct_box = slide.shapes.add_textbox(hero_left, hero_top + Inches(0.4), hero_width, Inches(0.7))
        pct_tf = pct_box.text_frame
        pct_tf.paragraphs[0].text = f"{best_plan.savings_pct:.0f}%"
        pct_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        pct_tf.paragraphs[0].font.name = 'Poppins'
        pct_tf.paragraphs[0].font.size = Pt(48)
        pct_tf.paragraphs[0].font.bold = True
        pct_tf.paragraphs[0].font.color.rgb = text_color

        # Add "annually" or savings amount subtitle
        subtitle_box = slide.shapes.add_textbox(hero_left, hero_top + Inches(1.05), hero_width, Inches(0.3))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.paragraphs[0].text = f"({self._format_currency(abs(best_plan.savings_amount))}/year)"
        subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_tf.paragraphs[0].font.name = 'Poppins'
        subtitle_tf.paragraphs[0].font.size = Pt(11)
        subtitle_tf.paragraphs[0].font.color.rgb = text_color

    def _add_decoratives(self, slide) -> None:
        """Add decorative edge images to the slide"""
        # Edge image - bottom right (smaller for better balance)
        edge_width = Inches(3.2)
        edge_height = Inches(0.8)
        if EDGE_IMAGE.exists():
            slide.shapes.add_picture(
                str(EDGE_IMAGE),
                left=SLIDE_WIDTH - edge_width,
                top=SLIDE_HEIGHT - edge_height,
                width=edge_width,
                height=edge_height
            )

    def generate(self) -> BytesIO:
        """
        Generate PowerPoint and return as BytesIO buffer.

        Returns:
            BytesIO buffer containing the .pptx file
        """
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer

    def save(self, output_path: str) -> None:
        """
        Save PowerPoint to file.

        Args:
            output_path: Path to save the file
        """
        self.prs.save(output_path)


def generate_cooperative_health_slide(data: CooperativeHealthData) -> BytesIO:
    """
    Convenience function to generate a Cooperative Health slide.

    Args:
        data: CooperativeHealthData with values to populate

    Returns:
        BytesIO buffer with generated PowerPoint
    """
    generator = CooperativeHealthSlideGenerator()
    generator.create_slide(data)
    return generator.generate()


def generate_cooperative_health_from_session(session_state) -> BytesIO:
    """
    Generate Cooperative Health slide from Streamlit session state.

    Args:
        session_state: st.session_state object

    Returns:
        BytesIO buffer with generated PowerPoint
    """
    data = CooperativeHealthData.from_session_state(session_state)
    return generate_cooperative_health_slide(data)


if __name__ == "__main__":
    # Test with sample data
    print("Testing Cooperative Health Slide Generator...")

    # Create sample data
    data = CooperativeHealthData(
        has_gap=True,
        total_gap_monthly=2464,
        gap_rate_ee=119,
        gap_rate_es=212,
        gap_rate_ec=187,
        gap_rate_f=245,
        tiers=[
            TierData(
                name="Employee Only",
                code="EE",
                current_total=4523,
                current_base=4047,
                current_gap=476,
                renewal_total=4978,
                renewal_base=4502,
                renewal_gap=476,
                hap_1k_total=3845,
                hap_1k_min=303,
                hap_1k_max=565,
                hap_25k_total=3412,
                hap_25k_min=268,
                hap_25k_max=498,
            ),
            TierData(
                name="Employee + Spouse",
                code="ES",
                current_total=3245,
                current_base=2821,
                current_gap=424,
                renewal_total=3570,
                renewal_base=3146,
                renewal_gap=424,
                hap_1k_total=2756,
                hap_1k_min=458,
                hap_1k_max=847,
                hap_25k_total=2445,
                hap_25k_min=405,
                hap_25k_max=748,
            ),
            TierData(
                name="Employee + Children",
                code="EC",
                current_total=2156,
                current_base=1782,
                current_gap=374,
                renewal_total=2372,
                renewal_base=1998,
                renewal_gap=374,
                hap_1k_total=1834,
                hap_1k_min=389,
                hap_1k_max=623,
                hap_25k_total=1627,
                hap_25k_min=344,
                hap_25k_max=551,
            ),
            TierData(
                name="Family",
                code="F",
                current_total=2308,
                current_base=1818,
                current_gap=490,
                renewal_total=2539,
                renewal_base=2049,
                renewal_gap=490,
                hap_1k_total=1839,
                hap_1k_min=534,
                hap_1k_max=892,
                hap_25k_total=1632,
                hap_25k_min=472,
                hap_25k_max=788,
            ),
        ],
        total_current=12232,
        total_renewal=13459,
        total_hap_1k=10274,
        total_hap_25k=9116,
        client_name="Sample Company",
    )

    # Generate slide
    generator = CooperativeHealthSlideGenerator()
    generator.create_slide(data)

    # Save to test file
    output_path = Path(__file__).parent / 'test_cooperative_health.pptx'
    generator.save(str(output_path))

    print(f"Generated test slide: {output_path}")
    print("Test complete!")
