"""
Marketplace Rates by Coverage Type Slide Generator

Generates a single PowerPoint slide with the Marketplace Rates table
that matches the design from the ICHRA dashboard.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from typing import Dict, List, Optional
from io import BytesIO
from pathlib import Path
from dataclasses import dataclass, field
from datetime import datetime
import pandas as pd


# Color scheme matching ICHRA dashboard design
COLORS = {
    # Bronze
    'bronze_header_bg': RGBColor(0xFE, 0xF3, 0xE2),   # #FEF3E2
    'bronze_text': RGBColor(0x92, 0x40, 0x0E),        # #92400E
    'bronze_light': RGBColor(0xB4, 0x53, 0x09),       # #B45309

    # Silver
    'silver_header_bg': RGBColor(0xF3, 0xF4, 0xF6),   # #F3F4F6
    'silver_text': RGBColor(0x37, 0x41, 0x51),        # #374151
    'silver_light': RGBColor(0x6B, 0x72, 0x80),       # #6B7280

    # Gold
    'gold_header_bg': RGBColor(0xFE, 0xF9, 0xC3),     # #FEF9C3
    'gold_text': RGBColor(0x85, 0x4D, 0x0E),          # #854D0E
    'gold_light': RGBColor(0xA1, 0x62, 0x07),         # #A16207

    # Neutral
    'row_label_bg': RGBColor(0xF9, 0xFA, 0xFB),       # #F9FAFB
    'row_label_text': RGBColor(0x10, 0x18, 0x28),     # #101828
    'secondary_text': RGBColor(0x6A, 0x72, 0x82),     # #6A7282
    'total_bg': RGBColor(0xF3, 0xF4, 0xF6),           # #F3F4F6
    'total_text': RGBColor(0x10, 0x18, 0x28),         # #101828
    'border': RGBColor(0xE5, 0xE7, 0xEB),             # #E5E7EB
    'white': RGBColor(0xFF, 0xFF, 0xFF),              # #FFFFFF

    # Savings/highlight
    'savings_green': RGBColor(0x00, 0xA6, 0x3E),      # #00A63E
    'savings_bg': RGBColor(0xDC, 0xFC, 0xE7),         # #DCFCE7
    'cost_red': RGBColor(0xDC, 0x26, 0x26),           # #DC2626

    # Renewal comparison row (more visible gray)
    'renewal_bg': RGBColor(0xE5, 0xE7, 0xEB),         # #E5E7EB (gray-200)
    'renewal_text': RGBColor(0x37, 0x41, 0x51),       # #374151 (gray-700)
}

# Slide dimensions (standard widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Decorative image paths
DECORATIVE_IMAGES_DIR = Path(__file__).parent / "glove-design" / "Design + Crit 01"
CORNER_IMAGE = DECORATIVE_IMAGES_DIR / "glove-tile-corner.png"
EDGE_IMAGE = DECORATIVE_IMAGES_DIR / "glove-tile-edge-h-fade.png"

# Banner image path
BANNER_IMAGE = Path("/Users/jimdonovan/Desktop/GLOVE/PPT_header.png")


@dataclass
class TierMetalData:
    """Data for a single metal level within a tier"""
    min_rate: float = 0.0
    max_rate: float = 0.0
    total: float = 0.0


@dataclass
class TierData:
    """Data for a single coverage tier row"""
    name: str  # Display name (e.g., "Employee Only")
    code: str  # Tier code (e.g., "EE")
    count: int = 0
    avg_age: int = 0
    bronze: TierMetalData = field(default_factory=TierMetalData)
    silver: TierMetalData = field(default_factory=TierMetalData)
    gold: TierMetalData = field(default_factory=TierMetalData)


@dataclass
class MarketplaceRatesData:
    """Data container for Marketplace Rates slide"""

    # Tier breakdown
    tiers: List[TierData] = field(default_factory=list)

    # Plan counts
    bronze_plan_count: int = 0
    silver_plan_count: int = 0
    gold_plan_count: int = 0

    # Metal AV percentages
    bronze_av: int = 60
    silver_av: int = 70
    gold_av: int = 80

    # Monthly totals
    bronze_monthly: float = 0.0
    silver_monthly: float = 0.0
    gold_monthly: float = 0.0

    # Annual totals
    bronze_annual: float = 0.0
    silver_annual: float = 0.0
    gold_annual: float = 0.0

    # Renewal premium for comparison
    renewal_monthly: float = 0.0

    # Client info
    client_name: str = ""
    generated_date: str = ""

    @classmethod
    def from_dashboard_data(cls, tier_costs: Dict, renewal_monthly: float = 0,
                            metal_av: Dict = None, client_name: str = "") -> 'MarketplaceRatesData':
        """
        Build MarketplaceRatesData from calculate_tier_marketplace_costs() output.

        Args:
            tier_costs: Output from calculate_tier_marketplace_costs()
            renewal_monthly: Current renewal premium for savings calculation
            metal_av: Dict of metal level -> AV percentage
            client_name: Client/company name

        Returns:
            Populated MarketplaceRatesData instance
        """
        data = cls()
        data.client_name = client_name
        data.generated_date = datetime.now().strftime("%m.%d.%y")
        data.renewal_monthly = renewal_monthly

        if metal_av:
            data.bronze_av = int(metal_av.get('Bronze', 60))
            data.silver_av = int(metal_av.get('Silver', 70))
            data.gold_av = int(metal_av.get('Gold', 80))

        # Plan counts
        plan_counts = tier_costs.get('plan_counts', {})
        data.bronze_plan_count = plan_counts.get('Bronze', 0)
        data.silver_plan_count = plan_counts.get('Silver', 0)
        data.gold_plan_count = plan_counts.get('Gold', 0)

        # Totals
        totals = tier_costs.get('totals', {})
        data.bronze_monthly = totals.get('Bronze', {}).get('monthly', 0)
        data.silver_monthly = totals.get('Silver', {}).get('monthly', 0)
        data.gold_monthly = totals.get('Gold', {}).get('monthly', 0)
        data.bronze_annual = totals.get('Bronze', {}).get('annual', 0)
        data.silver_annual = totals.get('Silver', {}).get('annual', 0)
        data.gold_annual = totals.get('Gold', {}).get('annual', 0)

        # Tier data
        tiers_data = tier_costs.get('tiers', {})
        tier_order = ['Employee Only', 'Employee + Spouse', 'Employee + Children', 'Family']
        tier_codes = {'Employee Only': 'EE', 'Employee + Spouse': 'ES',
                      'Employee + Children': 'EC', 'Family': 'F'}

        for tier_name in tier_order:
            tier_info = tiers_data.get(tier_name)
            if not tier_info:
                continue

            count = tier_info.get('count', 0)
            if count == 0:
                continue

            tier = TierData(
                name=tier_name,
                code=tier_codes.get(tier_name, ''),
                count=count,
                avg_age=tier_info.get('avg_age', 0),
            )

            # Metal data
            for metal, attr in [('Bronze', 'bronze'), ('Silver', 'silver'), ('Gold', 'gold')]:
                metal_info = tier_info.get(metal, {})
                setattr(tier, attr, TierMetalData(
                    min_rate=metal_info.get('min', 0),
                    max_rate=metal_info.get('max', 0),
                    total=metal_info.get('total', 0),
                ))

            data.tiers.append(tier)

        return data


class MarketplaceRatesSlideGenerator:
    """Generate Marketplace Rates by Coverage Type PowerPoint slide"""

    def __init__(self):
        """Initialize with a new blank presentation"""
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    def _set_cell_fill(self, cell, color: RGBColor):
        """Set cell background color"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    def _set_cell_text(self, cell, text: str, color: RGBColor,
                       font_size: int = 11, bold: bool = False,
                       align: PP_ALIGN = PP_ALIGN.CENTER, italic: bool = False):
        """Set cell text with formatting"""
        cell.text = text
        cell.text_frame.paragraphs[0].alignment = align
        cell.text_frame.paragraphs[0].font.name = 'Poppins'
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.text_frame.paragraphs[0].font.color.rgb = color
        cell.text_frame.paragraphs[0].font.bold = bold
        cell.text_frame.paragraphs[0].font.italic = italic
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _add_cell_line(self, cell, text: str, color: RGBColor,
                       font_size: int = 10, bold: bool = False,
                       align: PP_ALIGN = PP_ALIGN.CENTER):
        """Add additional line to cell"""
        p = cell.text_frame.add_paragraph()
        p.text = text
        p.alignment = align
        p.font.name = 'Poppins'
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold

    def _format_currency(self, value: float) -> str:
        """Format value as currency"""
        if value and value > 0:
            return f"${value:,.0f}"
        return "--"

    def _format_range(self, min_val: float, max_val: float) -> str:
        """Format rate range"""
        if min_val > 0 and max_val > 0:
            if min_val == max_val:
                return f"${min_val:,.0f}"
            return f"${min_val:,.0f}-${max_val:,.0f}"
        return "--"

    def _calc_savings(self, metal_monthly: float, renewal_monthly: float) -> tuple:
        """Calculate savings vs renewal. Returns (text, color)"""
        if renewal_monthly <= 0 or metal_monthly <= 0:
            return ("--", COLORS['secondary_text'])
        diff = renewal_monthly - metal_monthly
        if diff > 0:
            return (f"-${diff:,.0f}", COLORS['savings_green'])
        elif diff < 0:
            return (f"${abs(diff):,.0f}", COLORS['cost_red'])  # No sign needed, red color indicates cost increase
        return ("--", COLORS['secondary_text'])

    def _calc_savings_pct(self, metal_monthly: float, renewal_monthly: float) -> tuple:
        """Calculate savings percentage. Returns (text, color)"""
        if renewal_monthly <= 0 or metal_monthly <= 0:
            return ("--", COLORS['secondary_text'])
        diff = renewal_monthly - metal_monthly
        pct = (diff / renewal_monthly) * 100
        if pct > 0:
            return (f"{pct:.0f}%", COLORS['savings_green'])
        return ("--", COLORS['secondary_text'])

    def create_slide(self, data: MarketplaceRatesData) -> None:
        """
        Create the Marketplace Rates slide.

        Args:
            data: MarketplaceRatesData with values to display
        """
        # Add blank slide
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # Add banner at top
        if BANNER_IMAGE.exists():
            slide.shapes.add_picture(
                str(BANNER_IMAGE),
                left=Inches(0),
                top=Inches(0),
                width=SLIDE_WIDTH,
                height=Inches(0.25)
            )

        # Banner offset for content positioning
        banner_offset = Inches(0.2)

        # Add company name (top right)
        if data.client_name:
            company_box = slide.shapes.add_textbox(
                Inches(9.5), Inches(0.3), Inches(3.5), Inches(0.3)
            )
            company_tf = company_box.text_frame
            company_tf.paragraphs[0].text = data.client_name
            company_tf.paragraphs[0].font.name = 'Poppins'
            company_tf.paragraphs[0].font.size = Pt(14)
            company_tf.paragraphs[0].font.bold = True
            company_tf.paragraphs[0].font.color.rgb = COLORS['total_text']
            company_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Add title
        title_left = Inches(0.5)
        title_top = Inches(0.35) + banner_offset
        title_width = Inches(9)
        title_height = Inches(0.6)

        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_tf = title_box.text_frame
        title_tf.paragraphs[0].text = "Marketplace Rates by Coverage Type"
        title_tf.paragraphs[0].font.name = 'Poppins'
        title_tf.paragraphs[0].font.size = Pt(26)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = COLORS['total_text']

        # Add subtitle
        subtitle_top = Inches(0.85) + banner_offset
        subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top, Inches(12.33), Inches(0.4))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.paragraphs[0].text = "Lowest cost plans by metal level"
        subtitle_tf.paragraphs[0].font.name = 'Poppins'
        subtitle_tf.paragraphs[0].font.size = Pt(14)
        subtitle_tf.paragraphs[0].font.color.rgb = COLORS['secondary_text']

        # Calculate rows: header + plans available + tiers + total + annual + savings + savings %
        num_tiers = len(data.tiers)
        num_rows = 1 + 1 + num_tiers + 4  # header, plans available, tiers, footer rows
        num_cols = 4  # Coverage Type, Bronze, Silver, Gold

        # Create table
        table_left = Inches(0.5)
        table_top = Inches(1.4) + banner_offset
        table_width = Inches(12.33)
        table_height = Inches(0.5 + 0.4 + (num_tiers * 0.6) + (4 * 0.45))

        table = slide.shapes.add_table(num_rows, num_cols, table_left, table_top,
                                        table_width, table_height).table

        # Set column widths
        table.columns[0].width = Inches(3.5)
        for i in range(1, 4):
            table.columns[i].width = Inches(2.94)

        # Row 0: Header
        header_configs = [
            ("Coverage Type", COLORS['white'], COLORS['row_label_text']),
            (f"Bronze\n{data.bronze_av}% AV", COLORS['bronze_header_bg'], COLORS['bronze_text']),
            (f"Silver\n{data.silver_av}% AV", COLORS['silver_header_bg'], COLORS['silver_text']),
            (f"Gold\n{data.gold_av}% AV", COLORS['gold_header_bg'], COLORS['gold_text']),
        ]

        for col_idx, (text, bg_color, text_color) in enumerate(header_configs):
            cell = table.cell(0, col_idx)
            self._set_cell_fill(cell, bg_color)
            # Handle multi-line headers
            lines = text.split('\n')
            self._set_cell_text(cell, lines[0], text_color, font_size=14, bold=True,
                               align=PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER)
            if len(lines) > 1:
                self._add_cell_line(cell, lines[1], text_color, font_size=10, bold=False)

        # Row 1: Plans available
        cell = table.cell(1, 0)
        self._set_cell_fill(cell, COLORS['row_label_bg'])
        self._set_cell_text(cell, "Plans available", COLORS['secondary_text'],
                           font_size=11, align=PP_ALIGN.LEFT)

        plan_counts = [data.bronze_plan_count, data.silver_plan_count, data.gold_plan_count]
        metal_text_colors = [COLORS['bronze_text'], COLORS['silver_text'], COLORS['gold_text']]

        for col_idx, (count, text_color) in enumerate(zip(plan_counts, metal_text_colors), start=1):
            cell = table.cell(1, col_idx)
            self._set_cell_fill(cell, COLORS['row_label_bg'])
            self._set_cell_text(cell, f"{count:,}" if count > 0 else "--", COLORS['secondary_text'], font_size=12)

        # Determine which metal column has lowest total (for consistent column highlighting)
        monthly_totals = {'Bronze': data.bronze_monthly, 'Silver': data.silver_monthly, 'Gold': data.gold_monthly}
        valid_totals = {k: v for k, v in monthly_totals.items() if v > 0}
        lowest_metal = min(valid_totals, key=valid_totals.get) if valid_totals else None
        lowest_col = {'Bronze': 1, 'Silver': 2, 'Gold': 3}.get(lowest_metal)

        # Tier rows with column-based highlighting
        for row_idx, tier in enumerate(data.tiers, start=2):
            # Coverage Type column
            cell = table.cell(row_idx, 0)
            self._set_cell_fill(cell, COLORS['row_label_bg'])
            self._set_cell_text(cell, tier.name, COLORS['row_label_text'],
                               font_size=12, bold=True, align=PP_ALIGN.LEFT)
            age_suffix = f" | avg age {tier.avg_age}" if tier.avg_age > 0 else ""
            self._add_cell_line(cell, f"{tier.count} employee{'s' if tier.count != 1 else ''}{age_suffix}",
                               COLORS['secondary_text'], font_size=10, align=PP_ALIGN.LEFT)

            # Metal columns - highlight entire column if lowest total
            for col_idx, (metal_data, text_color) in enumerate([
                (tier.bronze, COLORS['bronze_text']),
                (tier.silver, COLORS['silver_text']),
                (tier.gold, COLORS['gold_text']),
            ], start=1):
                cell = table.cell(row_idx, col_idx)
                self._set_cell_fill(cell, COLORS['white'])
                self._set_cell_text(cell, self._format_currency(metal_data.total),
                                   text_color, font_size=14, bold=True)
                self._add_cell_line(cell, self._format_range(metal_data.min_rate, metal_data.max_rate),
                                   COLORS['secondary_text'], font_size=10)

        # Footer rows
        footer_row_start = 2 + num_tiers

        # Total Monthly row
        row_idx = footer_row_start
        cell = table.cell(row_idx, 0)
        self._set_cell_fill(cell, COLORS['white'])
        self._set_cell_text(cell, "Total Monthly", COLORS['total_text'],
                           font_size=12, bold=True, align=PP_ALIGN.LEFT)

        monthly_values = [data.bronze_monthly, data.silver_monthly, data.gold_monthly]
        for col_idx, value in enumerate(monthly_values, start=1):
            cell = table.cell(row_idx, col_idx)
            self._set_cell_fill(cell, COLORS['white'])
            # Use dark text color for readability
            self._set_cell_text(cell, self._format_currency(value), COLORS['total_text'], font_size=14, bold=True)

        # "vs a renewal of $X/mo:" merged row
        row_idx = footer_row_start + 1
        # Merge all cells in this row
        start_cell = table.cell(row_idx, 0)
        end_cell = table.cell(row_idx, 3)
        start_cell.merge(end_cell)
        merged_cell = table.cell(row_idx, 0)
        self._set_cell_fill(merged_cell, COLORS['renewal_bg'])
        renewal_text = f"vs a renewal of {self._format_currency(data.renewal_monthly)}/mo:"
        self._set_cell_text(merged_cell, renewal_text, COLORS['renewal_text'],
                           font_size=14, bold=True, align=PP_ALIGN.CENTER, italic=True)

        # Monthly Savings row
        row_idx = footer_row_start + 2
        cell = table.cell(row_idx, 0)
        self._set_cell_fill(cell, COLORS['white'])
        self._set_cell_text(cell, "Monthly Savings", COLORS['total_text'],
                           font_size=12, bold=True, align=PP_ALIGN.LEFT)

        for col_idx, value in enumerate(monthly_values, start=1):
            cell = table.cell(row_idx, col_idx)
            self._set_cell_fill(cell, COLORS['white'])
            text, color = self._calc_savings(value, data.renewal_monthly)
            self._set_cell_text(cell, text, color, font_size=14, bold=True)

        # Premium savings % row
        row_idx = footer_row_start + 3
        cell = table.cell(row_idx, 0)
        self._set_cell_fill(cell, COLORS['white'])
        self._set_cell_text(cell, "Premium Savings %", COLORS['total_text'],
                           font_size=12, bold=True, align=PP_ALIGN.LEFT)

        for col_idx, value in enumerate(monthly_values, start=1):
            cell = table.cell(row_idx, col_idx)
            self._set_cell_fill(cell, COLORS['white'])
            text, color = self._calc_savings_pct(value, data.renewal_monthly)
            self._set_cell_text(cell, text, color, font_size=16, bold=True)

        # Add footer note
        note_top = table_top + table_height + Inches(0.2)
        note_box = slide.shapes.add_textbox(table_left, note_top, Inches(10), Inches(0.4))
        note_tf = note_box.text_frame
        note_tf.paragraphs[0].text = ("Rate ranges show lowest cost plan rates from youngest to oldest "
                                       "age band within each coverage tier.")
        note_tf.paragraphs[0].font.name = 'Poppins'
        note_tf.paragraphs[0].font.size = Pt(10)
        note_tf.paragraphs[0].font.color.rgb = COLORS['secondary_text']

        # Add generated footer
        footer_top = Inches(6.9)
        footer_box = slide.shapes.add_textbox(table_left, footer_top, Inches(10), Inches(0.3))
        footer_tf = footer_box.text_frame
        client_suffix = f" for {data.client_name}" if data.client_name else ""
        footer_tf.paragraphs[0].text = f"Generated by Glove Benefits{client_suffix} | {data.generated_date} | ICHRA Calculator"
        footer_tf.paragraphs[0].font.name = 'Poppins'
        footer_tf.paragraphs[0].font.size = Pt(9)
        footer_tf.paragraphs[0].font.color.rgb = COLORS['secondary_text']

        # Add decorative images
        self._add_decoratives(slide)

    def _add_decoratives(self, slide) -> None:
        """Add decorative corner and edge images to the slide"""
        corner_size = Inches(1.5)
        if CORNER_IMAGE.exists():
            slide.shapes.add_picture(
                str(CORNER_IMAGE),
                left=Inches(0),
                top=Inches(0),
                width=corner_size,
                height=corner_size
            )

        edge_width = Inches(3.2)
        edge_height = Inches(0.8)
        if EDGE_IMAGE.exists():
            slide.shapes.add_picture(
                str(EDGE_IMAGE),
                left=SLIDE_WIDTH - edge_width,
                top=SLIDE_HEIGHT - edge_height,
                width=edge_width,
                height=edge_height
            )

    def generate(self) -> BytesIO:
        """Generate PowerPoint and return as BytesIO buffer."""
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer

    def save(self, output_path: str) -> None:
        """Save PowerPoint to file."""
        self.prs.save(output_path)


def generate_marketplace_rates_slide(data: MarketplaceRatesData) -> BytesIO:
    """
    Convenience function to generate a Marketplace Rates slide.

    Args:
        data: MarketplaceRatesData with values to populate

    Returns:
        BytesIO buffer with generated PowerPoint
    """
    generator = MarketplaceRatesSlideGenerator()
    generator.create_slide(data)
    return generator.generate()


def generate_marketplace_rates_from_session(session_state) -> BytesIO:
    """
    Generate Marketplace Rates slide from Streamlit session state.

    Args:
        session_state: st.session_state object

    Returns:
        BytesIO buffer with generated PowerPoint
    """
    from pages import calculate_tier_marketplace_costs

    census_df = session_state.get('census_df')
    financial_summary = session_state.get('financial_summary', {})
    multi_metal_results = financial_summary.get('multi_metal_scenario', {})
    client_name = session_state.get('client_name', '')
    db = session_state.get('db')

    # Calculate tier costs
    tier_costs = calculate_tier_marketplace_costs(
        census_df,
        multi_metal_results=multi_metal_results,
        db=db
    )

    # Get renewal and metal AV
    renewal_monthly = financial_summary.get('renewal_2026_total', 0)

    data = MarketplaceRatesData.from_dashboard_data(
        tier_costs=tier_costs,
        renewal_monthly=renewal_monthly,
        metal_av={'Bronze': 60, 'Silver': 70, 'Gold': 80},
        client_name=client_name
    )

    return generate_marketplace_rates_slide(data)


if __name__ == "__main__":
    # Test with sample data
    print("Testing Marketplace Rates Slide Generator...")

    # Create sample data
    data = MarketplaceRatesData(
        bronze_plan_count=45,
        silver_plan_count=38,
        gold_plan_count=22,
        bronze_av=60,
        silver_av=70,
        gold_av=80,
        bronze_monthly=8500,
        silver_monthly=9200,
        gold_monthly=10800,
        bronze_annual=102000,
        silver_annual=110400,
        gold_annual=129600,
        renewal_monthly=12500,
        client_name="Sample Company",
        tiers=[
            TierData(
                name="Employee Only",
                code="EE",
                count=23,
                avg_age=38,
                bronze=TierMetalData(min_rate=303, max_rate=565, total=6969),
                silver=TierMetalData(min_rate=356, max_rate=645, total=7845),
                gold=TierMetalData(min_rate=412, max_rate=780, total=9234),
            ),
            TierData(
                name="Employee + Spouse",
                code="ES",
                count=8,
                avg_age=42,
                bronze=TierMetalData(min_rate=606, max_rate=1130, total=4848),
                silver=TierMetalData(min_rate=712, max_rate=1290, total=5696),
                gold=TierMetalData(min_rate=824, max_rate=1560, total=6592),
            ),
            TierData(
                name="Employee + Children",
                code="EC",
                count=5,
                avg_age=35,
                bronze=TierMetalData(min_rate=450, max_rate=780, total=2250),
                silver=TierMetalData(min_rate=530, max_rate=890, total=2650),
                gold=TierMetalData(min_rate=615, max_rate=1050, total=3075),
            ),
            TierData(
                name="Family",
                code="F",
                count=4,
                avg_age=40,
                bronze=TierMetalData(min_rate=850, max_rate=1450, total=3400),
                silver=TierMetalData(min_rate=998, max_rate=1680, total=3992),
                gold=TierMetalData(min_rate=1155, max_rate=1980, total=4620),
            ),
        ],
    )

    # Generate slide
    generator = MarketplaceRatesSlideGenerator()
    generator.create_slide(data)

    # Save to test file
    output_path = Path(__file__).parent / 'test_marketplace_rates.pptx'
    generator.save(str(output_path))

    print(f"Generated test slide: {output_path}")
    print("Test complete!")
